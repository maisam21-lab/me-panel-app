"""ME Sales Panel - Streamlit dashboard (BI-infographic style).

Single-canvas dashboard over the BigQuery bridge that also feeds the Google
Sheets panel (css-operations.me_panel_dev_us.me_sales_panel_k_monthly).
Sections: KPI gauges + cards / Revenue / Sales / Mix / Countries (animated map).

Credentials: [gcp_service_account] in .streamlit/secrets.toml, else ADC.
"""
import base64
import math
import os
from decimal import Decimal

import pandas as pd
import streamlit as st

BQ_PROJECT = "css-operations"
BRIDGE_TABLE = "css-operations.me_panel_dev_us.me_sales_panel_k_monthly"
COUNTRIES = ["Middle East", "Saudi Arabia", "UAE", "Kuwait", "Bahrain", "Qatar"]
# Historical reference only; load_bridge no longer filters by a lower bound so the app
# matches the Google Sheets panel's full-history view.
START_MONTH = "2023-08-31"   # panel history floor: August 2023 onward

# NAMAA brand palette: deep forest green + cream + terracotta (from the brand banner),
# with sage / gold / taupe as supporting series colors. Variable names kept generic so the
# chart code reads unchanged - only the values are branded.
NAVY = "#21362B"          # NAMAA deep green (primary dark)
TEAL = "#5F8575"          # sage green (secondary)
RED = "#B4472E"           # rust - losses/churn
ORANGE = "#D97757"        # terracotta - the brand accent
YELLOW = "#C2A14D"        # sand/gold
SLATE = "#A79E8B"         # taupe
LIGHT = "#EAE7DC"         # warm light (gauge remainder)

SERIES_COLORS = [NAVY, ORANGE, YELLOW, TEAL, RED, SLATE]
COUNTRY_COLORS = {"Saudi Arabia": NAVY, "UAE": ORANGE, "Kuwait": YELLOW, "Bahrain": RED, "Qatar": TEAL}

LOGO_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static", "namaa_logo.jpg")
_page_icon = LOGO_PATH if os.path.exists(LOGO_PATH) else ":bar_chart:"

st.set_page_config(page_title="ME Sales Panel | NAMAA", layout="wide", page_icon=_page_icon,
                   initial_sidebar_state="collapsed")


MOTIF_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static", "namaa_motif.png")

# ArcGIS Location Platform key for the Esri basemap tiles (same map service family as the
# Talabat site). Override via ARCGIS_API_KEY in secrets; this fallback keeps the map working.
ARCGIS_DEFAULT_KEY = ("AAPTazJyeyXdia4yZ8kifX4yDiw..tmKsfsr3ICm0GmhX-8bLyCDCls6oQi4tGrZizEYySs_"
                      "5IPHbtB5rkpdZc7WBoRiAp9PGgI0sp-WRrnzzuvrzbYqasb1qB_altV_pkh6jbu6YIUEk386"
                      "TZhBDRRN9AQ_nxQfLIE7bVHSdU7gKEABklA-fj-Oc_ktthlJvkHTkYRNwSy9R8Zna-rovVMy"
                      "lRxjlwuqX9AuEWbcA-8sJUlH-G8FandT3O1Dj14qi7_ntgKh5aC_YrK9ebg..AT1_mzOu9Cgs")


@st.cache_data(show_spinner=False)
def _logo_b64() -> str:
    try:
        with open(LOGO_PATH, "rb") as f:
            return base64.b64encode(f.read()).decode("ascii")
    except Exception:
        return ""


@st.cache_data(show_spinner=False)
def _motif_b64() -> str:
    try:
        with open(MOTIF_PATH, "rb") as f:
            return base64.b64encode(f.read()).decode("ascii")
    except Exception:
        return ""


def _inject_motif():
    """The terracotta NAMAA line-art, pinned to the top-right corner like the brand banner."""
    _mb = _motif_b64()
    if _mb:
        st.markdown('<img class="nm-motif" src="data:image/png;base64,' + _mb + '"/>',
                    unsafe_allow_html=True)


st.markdown(
    """
    <style>
    .stApp { background: #EEEDE5; }
    #MainMenu, footer { visibility: hidden; }
    .block-container { padding-top: 3.4rem !important; max-width: 100% !important;
                       padding-left: 1.6rem !important; padding-right: 1.6rem !important; }
    header[data-testid="stHeader"] { background: transparent !important; }
    /* NAMAA brand banner */
    .nm-banner { display: flex; align-items: center; gap: 16px; background: #21362B;
                 border-radius: 14px; padding: 10px 18px; margin: 4px 0 12px 0;
                 box-shadow: 0 2px 8px rgba(33, 54, 43, 0.25); }
    .nm-banner img { height: 56px; border-radius: 10px; }
    .nm-name { color: #FFFFFF; font-weight: 800; font-size: 1.1rem; letter-spacing: 0.3em; margin: 0; }
    .nm-sub { color: #C9D5CC; font-size: 0.74rem; font-weight: 700; letter-spacing: 0.14em;
              text-transform: uppercase; margin: 0; }
    .hdr { display: flex; align-items: baseline; gap: 14px; margin-bottom: 4px; }
    .hdr .t { font-size: 1.5rem; font-weight: 800; color: #21362B; letter-spacing: -0.02em; }
    .hdr .badge { background: #E2E6DC; color: #21362B; font-size: 0.75rem; font-weight: 700;
                  padding: 3px 10px; border-radius: 999px; }
    div[data-testid="stVerticalBlockBorderWrapper"] {
        background: #FFFFFF; border: 1px solid #E0DCCE !important; border-radius: 14px;
        box-shadow: 0 2px 6px rgba(33, 54, 43, 0.08);
        transition: box-shadow 0.15s ease, transform 0.15s ease;
    }
    div[data-testid="stVerticalBlockBorderWrapper"]:hover {
        box-shadow: 0 6px 16px rgba(33, 54, 43, 0.16); transform: translateY(-1px);
        position: relative; z-index: 60;
    }
    .kpi-l { font-size: 0.72rem; font-weight: 700; letter-spacing: 0.06em; text-transform: uppercase;
             color: #7C776A; margin: 0; }
    .kpi-v { font-size: 1.5rem; font-weight: 800; color: #21362B; margin: 2px 0 0 0; line-height: 1.1; }
    .kpi-d-up { color: #3F7A52; font-size: 0.78rem; font-weight: 700; }
    .kpi-d-dn { color: #B4472E; font-size: 0.78rem; font-weight: 700; }
    .kpi-d-na { color: #A79E8B; font-size: 0.78rem; font-weight: 600; }
    .g-lbl { text-align: center; font-size: 0.78rem; font-weight: 700; letter-spacing: 0.04em;
             text-transform: uppercase; color: #55604F; margin: -6px 0 0 0; }
    .g-dlt { text-align: center; font-size: 0.74rem; margin: 0; }
    .sec { font-size: 1rem; font-weight: 800; color: #21362B; text-transform: uppercase;
           letter-spacing: 0.06em; margin: 22px 0 4px 0;
           border-left: 5px double #D97757; padding-left: 10px; }
    section[data-testid="stSidebar"], [data-testid="collapsedControl"] { display: none !important; }
    .stDataFrame thead th { background: #F0EEE6 !important; font-weight: 600 !important; }
    /* NAMAA terracotta line-art, top-right corner (brand banner motif) */
    .nm-motif { position: fixed; top: 36px; right: 0; width: 320px; opacity: 0.9;
                z-index: 0; pointer-events: none; }
    /* Auto-insight line under section headers */
    .nm-insight { background: #FBFAF3; border-left: 3px solid #D97757; border-radius: 8px;
                  padding: 7px 12px; margin: 2px 0 10px 0; color: #4A5548;
                  font-size: 0.84rem; }
    /* Executive metric x market grid (mirrors the All Hands slide mental model) */
    .exec-wrap { background: #FFFFFF; border: 1px solid #E0DCCE; border-radius: 14px;
                 padding: 6px 10px; box-shadow: 0 2px 6px rgba(33,54,43,0.08); overflow: visible; }
    @media (max-width: 1100px) { .exec-wrap { overflow-x: auto; } }
    .exec-grid { width: 100%; border-collapse: collapse; }
    .exec-grid th { background: #F0EEE6; color: #21362B; font-size: 0.78rem; font-weight: 800;
                    text-transform: uppercase; letter-spacing: 0.05em; padding: 8px 12px;
                    text-align: right; }
    .exec-grid th:first-child { text-align: left; }
    .exec-grid td { padding: 7px 12px; font-size: 0.9rem; text-align: right;
                    border-bottom: 1px solid #F1EEE4; color: #21362B; }
    .exec-grid td:first-child { text-align: left; font-weight: 600; }
    .exec-grid tr.exec-band td { background: #21362B; color: #FFFFFF; font-size: 0.7rem;
                                 font-weight: 800; letter-spacing: 0.1em; text-transform: uppercase;
                                 padding: 5px 12px; border-bottom: none; }
    .exec-grid .ev { font-weight: 800; }
    .exec-grid .ed { font-size: 0.72rem; display: block; margin-top: 1px; }
    /* Hover info cards: calculation + live behavior appear when hovering any wrapped number */
    .mi-wrap { position: relative; display: inline-block; cursor: help;
               border-bottom: 1px dotted rgba(124, 119, 106, 0.55); }
    .mi-wrap .mi-tip { display: none; position: absolute; top: calc(100% + 8px); left: 0;
                       width: 320px; max-width: 76vw; background: #21362B; color: #E7E4D8;
                       border-radius: 10px; padding: 10px 13px; font-size: 0.76rem; font-weight: 400;
                       line-height: 1.5; text-align: left; text-transform: none;
                       letter-spacing: normal; white-space: normal; z-index: 10000;
                       box-shadow: 0 10px 26px rgba(33, 54, 43, 0.4); pointer-events: none; }
    .mi-wrap:hover .mi-tip { display: block; }
    .mi-tip .tt { display: block; font-weight: 800; font-size: 0.78rem; color: #FFFFFF;
                  text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 4px; }
    .mi-tip .tb { display: block; margin-top: 5px; padding-top: 5px;
                  border-top: 1px solid rgba(238, 237, 229, 0.25); }
    .mi-tip .tw { display: block; margin-top: 5px; color: #E8B39B; }
    .mi-tip b { color: #FFFFFF; }
    .mi-right .mi-tip { left: auto; right: 0; }
    .mi-up .mi-tip { top: auto; bottom: calc(100% + 8px); }
    div[data-testid="stMarkdownContainer"] { overflow: visible !important; }
    /* NAMAA Copilot */
    .ai-hero { display: flex; align-items: center; gap: 16px; border-radius: 14px;
               padding: 16px 20px; margin: 4px 0 12px 0; color: #EEEDE5;
               background: linear-gradient(120deg, #21362B 0%, #2F4A3B 55%, #21362B 100%);
               box-shadow: 0 4px 14px rgba(33, 54, 43, 0.3); }
    .ai-spark { font-size: 1.9rem; color: #D97757; line-height: 1; }
    .ai-t { margin: 0; font-size: 1.15rem; font-weight: 800; letter-spacing: 0.08em;
            text-transform: uppercase; color: #FFFFFF; }
    .ai-s { margin: 2px 0 0 0; font-size: 0.82rem; color: #C9D5CC; }
    .ai-badge { margin-left: auto; font-size: 0.7rem; font-weight: 700; letter-spacing: 0.06em;
                text-transform: uppercase; padding: 4px 10px; border-radius: 999px; white-space: nowrap; }
    .ai-badge.on { background: rgba(95, 133, 117, 0.35); color: #D6E4DA; border: 1px solid #5F8575; }
    .ai-badge.off { background: rgba(217, 119, 87, 0.18); color: #E8B39B; border: 1px solid #D97757; }
    .anom-row { display: flex; flex-wrap: wrap; gap: 6px; align-items: center; margin: 2px 0 8px 0; }
    .anom-t { font-size: 0.7rem; font-weight: 800; letter-spacing: 0.08em; text-transform: uppercase;
              color: #7C776A; margin-right: 4px; }
    .anom { font-size: 0.74rem; font-weight: 700; padding: 3px 10px; border-radius: 999px; }
    .anom.warm { background: #F7E8DC; color: #A85A3A; border: 1px solid #D97757; }
    .anom.hot { background: #F4DCD4; color: #8F3A24; border: 1px solid #B4472E; }
    </style>
    """,
    unsafe_allow_html=True,
)


# ---------------------------------------------------------------- data layer
def _bq_client():
    from google.cloud import bigquery

    try:
        info = dict(st.secrets.get("gcp_service_account", {}))
    except Exception:
        info = {}
    if info and info.get("private_key"):
        try:
            from google.oauth2 import service_account

            info.setdefault("token_uri", "https://oauth2.googleapis.com/token")
            info.setdefault("auth_uri", "https://accounts.google.com/o/oauth2/auth")
            creds = service_account.Credentials.from_service_account_info(
                info, scopes=["https://www.googleapis.com/auth/bigquery"]
            )
            job_project = info.get("project_id") or BQ_PROJECT
            return bigquery.Client(project=job_project, credentials=creds)
        except Exception:
            pass
    return bigquery.Client(project=BQ_PROJECT)


@st.cache_data(ttl=900, show_spinner="Loading panel data from BigQuery...")
def load_bridge() -> pd.DataFrame:
    """Pull the full ME sales panel bridge from BigQuery. We only bound the upper end
    (LAST_DAY of the current month) so partial future rows do not show; the lower bound
    is left off so the app returns as many months of history as the bridge holds,
    matching the Google Sheets panel view."""
    client = _bq_client()
    rows = None
    try:
        query = (
            "SELECT * FROM `" + BRIDGE_TABLE + "` "
            "WHERE month_end >= DATE '" + START_MONTH + "' "
            "AND month_end <= LAST_DAY(CURRENT_DATE(), MONTH) "
            "ORDER BY month_end, country"
        )
        rows = list(client.query(query).result())
    except Exception as query_err:
        try:
            table = client.get_table(BRIDGE_TABLE)
            rows = list(client.list_rows(table))
        except Exception:
            raise query_err
    out = []
    for row in rows:
        d = dict(row)
        for k, v in list(d.items()):
            if isinstance(v, Decimal):
                d[k] = float(v)
        out.append(d)
    df = pd.DataFrame(out)
    if not df.empty:
        df["month_end"] = pd.to_datetime(df["month_end"])
        end = pd.Timestamp.today().normalize() + pd.offsets.MonthEnd(0)
        df = df[(df["month_end"] >= pd.Timestamp(START_MONTH)) & (df["month_end"] <= end)]
        df = df.sort_values(["month_end", "country"]).reset_index(drop=True)
    return df


def fmt(v, kind):
    try:
        if v is None or (isinstance(v, float) and math.isnan(v)):
            return ""
        if kind == "usd":
            return f"${float(v):,.0f}"
        if kind == "usdk":
            x = float(v)
            if abs(x) >= 1_000_000:
                return f"${x / 1_000_000:,.2f}M"
            if abs(x) >= 10_000:
                return f"${x / 1_000:,.0f}k"
            return f"${x:,.0f}"
        if kind == "pct":
            return f"{float(v) * 100:.1f}%"
        if kind == "num1":
            return f"{float(v):,.1f}"
        return f"{float(v):,.0f}"
    except Exception:
        return "" if v is None else str(v)


# Metrics offered in the Countries section (label, column, format kind).
COMPARE_METRICS = [
    ("CWs (kitchens)", "cws", "num"),
    ("Approved Deals", "approved_deals", "num"),
    ("New Occupied Kitchens", "new_occupied_k", "num"),
    ("NRRA $ (net)", "nrra_usd", "usd"),
    ("RRA $ (recognized)", "rra_usd", "usd"),
    ("RRL $ (recognized)", "rrl_usd", "usd"),
    ("RRX $ (accessed)", "xrra_usd", "usd"),
    ("Gross Recurring Revenue $ (EoP)", "gross_rr_usd", "usd"),
    ("RR after MKO/MFO $ (EoP)", "rr_after_mko_mfo_usd", "usd"),
    ("TCV $", "tcv_usd", "usd"),
    ("Approved TCV $", "approved_tcv_usd", "usd"),
    ("Churns (excl. transfers)", "churns_excl_transfers", "num"),
    ("Net Adds", "net_adds", "num"),
    ("Occupancy", "occupancy", "pct"),
    ("Occupied Kitchens", "occupied_kitchens", "num"),
    ("Live - Sold Rate %", "live_sold_rate", "pct"),
    ("Live - Sold Rate w/ Approved %", "live_sold_rate_approved", "pct"),
    ("Sold Rate - All Facilities", "sold_rate_all", "pct"),
]

COUNTRY_ISO = {"Saudi Arabia": "SAU", "UAE": "ARE", "Kuwait": "KWT", "Bahrain": "BHR", "Qatar": "QAT"}
COUNTRY_CENTROID = {  # (lat, lon) for the map bubbles; Bahrain nudged off Qatar
    "Saudi Arabia": (24.0, 44.5), "UAE": (23.6, 54.2), "Kuwait": (29.6, 47.6),
    "Bahrain": (26.9, 50.4), "Qatar": (24.9, 51.2),
}


# ---------------------------------------------------------------- access gate
def _allowed_emails() -> set:
    try:
        raw = st.secrets.get("ALLOWED_EMAILS", [])
    except Exception:
        return set()
    if isinstance(raw, str):
        raw = raw.replace(";", ",").split(",")
    return {str(e).strip().lower() for e in raw if str(e).strip()}


DEV_EMAILS_DEFAULT = {"maysam.abukashabeh@namaame.com"}


def _is_developer() -> bool:
    """Developer-only controls (e.g. the Refresh button). Override list via
    DEVELOPER_EMAILS in secrets. With no allowlist configured (local run) everyone
    on the machine is the developer."""
    try:
        raw = st.secrets.get("DEVELOPER_EMAILS", [])
    except Exception:
        raw = []
    if isinstance(raw, str):
        raw = raw.replace(";", ",").split(",")
    devs = {str(e).strip().lower() for e in raw if str(e).strip()} or DEV_EMAILS_DEFAULT
    if not _allowed_emails():
        return True
    return (st.session_state.get("me_user_email") or "").strip().lower() in devs


def _access_gate():
    allowed = _allowed_emails()
    if not allowed:
        return
    current = (st.session_state.get("me_user_email") or "").strip().lower()
    if current in allowed:
        return
    # Auto-identify allowed users - no typing:
    # 0) CSS data-apps (internal platform): Okta SSO at the platform level; streamlit_utils
    #    is only importable inside that environment and returns the verified Okta email.
    #    Same code elsewhere just hits ImportError and falls through.
    try:
        from streamlit_utils import auth as _css_auth
        _u = str(_css_auth.get_user_email() or "").strip()
        if _u and _u.lower() in allowed:
            st.session_state["me_user_email"] = _u
            return
    except Exception:
        pass
    # 1) Streamlit-authenticated viewer (when the platform knows who you are).
    try:
        _u = str(getattr(getattr(st, "user", None), "email", "") or "").strip()
        if _u and _u.lower() in allowed:
            st.session_state["me_user_email"] = _u
            return
    except Exception:
        pass
    # 2) Personalized link: https://<app-url>/?u=name@namaame.com
    #    Share each person THEIR link once; after that they never see a prompt.
    try:
        _qp = st.query_params.get("u", "")
        if isinstance(_qp, list):
            _qp = _qp[0] if _qp else ""
        _qp = str(_qp).strip()
        if _qp and _qp.lower() in allowed:
            st.session_state["me_user_email"] = _qp
            return
    except Exception:
        pass
    _b64 = _logo_b64()
    if _b64:
        st.markdown(
            '<div class="nm-banner"><img src="data:image/jpeg;base64,' + _b64 + '"/>'
            '<div><p class="nm-name">NAMAA</p><p class="nm-sub">ME Sales Panel</p></div></div>',
            unsafe_allow_html=True,
        )
    _inject_motif()
    st.write("Enter your work email to open the panel.")
    email = st.text_input("Email", key="me_email_input")
    if st.button("Open", type="primary"):
        if email.strip().lower() in allowed:
            st.session_state["me_user_email"] = email.strip()
            try:
                # Stamp identity into the URL so refreshes and bookmarks stay signed in
                # (the ?u= param re-passes the gate automatically on the next session).
                st.query_params["u"] = email.strip()
            except Exception:
                pass
            st.rerun()
        else:
            st.error("This email is not on the allowed list. Ask Maysam to add you.")
    st.stop()


# ---------------------------------------------------------------- chart kit
def _go():
    import plotly.graph_objects as go
    return go


def _base_layout(fig, title, height=340):
    fig.update_layout(
        title=dict(text=title, font=dict(size=14, color="#21362B", family="Arial Black, Arial")),
        height=height,
        margin=dict(l=8, r=8, t=44, b=4),
        legend=dict(orientation="h", y=-0.22, font=dict(size=11)),
        plot_bgcolor="white", paper_bgcolor="white",
        hovermode="x unified",
        font=dict(family="Arial", size=11, color="#334155"),
    )
    fig.update_xaxes(showgrid=False)
    fig.update_yaxes(gridcolor="#E7E4D8", griddash="dot", zerolinecolor="#D9D5C5")
    return fig


def money_axis(fig):
    fig.update_yaxes(tickprefix="$", tickformat="~s")
    return fig


def pct_axis(fig):
    fig.update_yaxes(tickformat=".0%")
    return fig


def add_line(fig, x, y, name, color, closed_idx, width=2.6, msize=8, fmt_kind=None):
    """Dot-marker line (reference style); the stretch after the last closed month is dotted.
    fmt_kind formats the hover value ($ / % / count) instead of raw floats."""
    go = _go()
    solid_end = min(closed_idx, len(x) - 1)
    cd = [fmt(v, fmt_kind) for v in y] if fmt_kind else None
    ht = "<b>%{fullData.name}</b>: %{customdata}<extra></extra>" if fmt_kind else None
    fig.add_trace(go.Scatter(
        x=x[: solid_end + 1], y=y[: solid_end + 1], mode="lines+markers", name=name,
        line=dict(color=color, width=width),
        marker=dict(size=msize, color="white", line=dict(color=color, width=2.5)),
        customdata=(cd[: solid_end + 1] if cd else None), hovertemplate=ht,
    ))
    if solid_end < len(x) - 1:
        fig.add_trace(go.Scatter(
            x=x[solid_end:], y=y[solid_end:], mode="lines+markers", name=name,
            line=dict(color=color, width=2, dash="dot"),
            marker=dict(size=msize - 2, color="white", line=dict(color=color, width=2)),
            opacity=0.55, showlegend=False, hoverinfo="skip",
        ))


def add_bar(fig, x, y, name, color, closed_idx, fmt_kind=None):
    go = _go()
    colors = [color if i <= closed_idx else "rgba(148,163,184,0.45)" for i in range(len(x))]
    cd = [fmt(v, fmt_kind) for v in y] if fmt_kind else None
    ht = "<b>%{fullData.name}</b>: %{customdata}<extra></extra>" if fmt_kind else None
    fig.add_trace(go.Bar(x=x, y=y, name=name, marker_color=colors,
                         customdata=cd, hovertemplate=ht))


def sec(title, subtitle=None):
    """Section header with an optional inline explainer."""
    sub = ('<span style="font-weight:600; font-size:0.78rem; color:#7C776A; '
           'text-transform:none; letter-spacing:0; margin-left:10px;">' + subtitle + "</span>"
           ) if subtitle else ""
    st.markdown(f'<p class="sec">{title}{sub}</p>', unsafe_allow_html=True)


def insight(text_html):
    """Auto-computed takeaway line under a section header."""
    st.markdown(f'<div class="nm-insight">{text_html}</div>', unsafe_allow_html=True)


# ---------------------------------------------------------------- metric explainers
# How every number is calculated - the honest bridge definitions (same logic that feeds
# the Google Sheets panel). (calculation, what-to-watch) per bridge column.
METRIC_INFO = {
    "cws": ("Contract Wins: distinct kitchen-level deals Closed-Won in the month (one deal = one "
            "kitchen), excluding member transfers. From Salesforce opportunities by Closed-Won date.",
            "The global-mart country definition excludes delayed transfers; facility views include them."),
    "approved_deals": ("Deals in Approved or Closed-Won stage counted by their approval date - built to "
                       "match the Salesforce Approved Deals report cell-by-cell (excludes Virtual/"
                       "CloudRetail kitchen types and member transfers).", ""),
    "cw_duration": ("Average contract length (months) of the month's Closed-Won deals, "
                    "LF-weighted.", ""),
    "new_occupied_k": ("Kitchens whose access date falls in the month, excluding member transfers. "
                       "Same access date the RRX $ row uses.",
                       "RRX $ includes transfer move-ins; this count does not - so on transfer-heavy "
                       "months $ can move without the count."),
    "rra": ("Recognized Recurring Revenue Added as % of the prior month's gross LF base.",
            "Recognized revenue matures over ~2 months; the live month shows a deal-date fill until "
            "recognized loads, then switches automatically."),
    "rrl": ("Recognized Recurring Revenue Lost as % of the prior month's gross LF base. License-Fee "
            "revenue only.", "Scheduled churns are known ahead - future months can show values."),
    "nrra": ("NRRA % = RRA % minus RRL %.", ""),
    "rra_usd": ("Recognized license-fee revenue of the month's Closed Wons (finance basis - same as "
                "the global panel).", "Live month = deal-date LF until recognized loads (~days after "
                "month close); numbers can then shift, e.g. Jun 2026 NRRA 111k->83k."),
    "rrl_usd": ("Recognized license-fee revenue lost to churn (finance basis).",
                "Backdated churns land retroactively - closed months converge over ~2 months."),
    "nrra_usd": ("NRRA $ = RRA $ minus RRL $.", "Inherits both parents' recognition timing."),
    "xrra_usd": ("RRX $: license fees of customers whose ACCESS date falls in the month (gross, "
                 "deal-date basis - no recognition lag). Transfers included.", ""),
    "xrrl_usd": ("RRLX $: license fees lost from post-access churns by churn date (pre-access churns "
                 "excluded).", ""),
    "nrrx_usd": ("NRRX $ = RRX $ minus RRLX $.", ""),
    "gross_rr_usd": ("Gross Recurring Revenue: every occupied kitchen's monthly license fee at end of "
                     "period. Occupied = latest accessed-and-not-churned Closed-Won deal per kitchen; "
                     "BP facilities included (reconciles with Occupied Kitchens).",
                     "Deals billed in non-local currency convert at their own currency's rate "
                     "(e.g. the Delivery Hero GBP kitchens)."),
    "rr_after_mko_mfo_usd": ("Same occupied-kitchen stock valued at the month's NET fee from the "
                             "Salesforce revenue schedules: LF minus MKO/MFO, custom, promo and term "
                             "discounts, with first/last-month proration (Total_MLF basis).",
                             "Gap vs Gross RR = the total concession load (~14% of gross for ME)."),
    "tcv_usd": ("Total Contract Value of the month's Closed Wons: monthly LF x contract length "
                "(x fx at the CW month).", ""),
    "approved_tcv_usd": ("TCV of Approved deals by approval month: monthly LF x contract length.", ""),
    "churns_excl_transfers": ("Kitchens churned in the month, transfers netted out (global-mart "
                              "definition).", ""),
    "net_adds": ("Net kitchen adds = CWs minus churns (both excl. transfers).", ""),
    "occupancy": ("Occupied kitchens (global mart) divided by Total Kitchen Numbers - TKN = the "
                  "account-declared kitchen count summed over LIVE facilities.", ""),
    "occupied_kitchens": ("Occupied kitchens from the global mart (active closed-won occupant "
                          "deals).", ""),
    "total_kitchens": ("TKN: SUM(account total_kitchen_numbers) over live facilities - the shared "
                       "denominator for Occupancy and the Live sold rates.", ""),
    "live_sold_rate": ("(Sold + Occupied + Churning kitchen statuses at LIVE facilities, BP included) "
                       "divided by TKN. Statuses are today's Salesforce snapshot.",
                       "SF keeps no kitchen-status history, so deep history is approximate; always >= "
                       "Occupancy by construction."),
    "live_sold_rate_approved": ("Live Sold Rate plus vacant kitchens carrying a still-pending Approved "
                                "deal (point-in-time pipeline: counted from approval until the deal "
                                "wins or dies).", "A pipeline stock - it does not tie to the monthly "
                                "Approved Deals report count."),
    "live_vacant_appr_k": ("Vacant kitchens at live facilities with a still-pending Approved deal as "
                           "of month-end (cumulative carry, drops when the deal closes).", ""),
    "sold_rate_all": ("Sold kitchens / all kitchens across all facilities (live + future).", ""),
    "net_sold_approved_rate": ("(Net sold + open approved pipeline) / all-facilities kitchens - same "
                               "denominator as Sold Rate All; the gap is the approved pipeline.", ""),
    "all_facilities": ("Count of owned sites (live + non-live facilities).", ""),
    "sales_team_size": ("In-seat AEs (weighted FTEs) - the Jad-locked productivity denominator.", ""),
    "sales_team_cw_productivity": ("CWs divided by in-seat sales team size.", ""),
    "aes": ("Account Executives (weighted headcount).", ""),
    "sdrs": ("SDRs (weighted headcount).", ""),
}


def _behavior_html(dd, col, kind, closed_idx):
    """'How this number is behaving': latest, MoM, 3-mo average, window high/low."""
    try:
        s = pd.to_numeric(dd[col].iloc[: closed_idx + 1], errors="coerce")
        labs = dd["month_label"].iloc[: closed_idx + 1]
        mask = s.notna()
        s, labs = s[mask].reset_index(drop=True), labs[mask].reset_index(drop=True)
        if s.empty:
            return ""
        parts = [f"latest <b>{fmt(s.iloc[-1], kind)}</b> ({labs.iloc[-1]})"]
        if len(s) >= 2:
            dv = s.iloc[-1] - s.iloc[-2]
            ar = "&#9650;" if dv >= 0 else "&#9660;"
            dtxt = f"{dv * 100:+.1f} pp" if kind == "pct" else ("+" if dv >= 0 else "") + fmt(dv, kind)
            parts.append(f"MoM {ar} {dtxt}")
        if len(s) >= 3:
            parts.append(f"3-mo avg {fmt(s.iloc[-3:].mean(), kind)}")
        pmax, pmin = int(s.values.argmax()), int(s.values.argmin())
        parts.append(f"high {fmt(s.iloc[pmax], kind)} ({labs.iloc[pmax]})")
        parts.append(f"low {fmt(s.iloc[pmin], kind)} ({labs.iloc[pmin]})")
        return " &middot; ".join(parts)
    except Exception:
        return ""


def _metric_tip(label, col, kind, dd, closed_idx, scope=None):
    """The hover info card body (.mi-tip span): calculation + live behavior + watch."""
    calc, watch = METRIC_INFO.get(col, ("", ""))
    beh = _behavior_html(dd, col, kind, closed_idx) if dd is not None else ""
    if not (calc or beh):
        return ""
    head = label + ((" &middot; " + scope) if scope else "")
    inner = f'<span class="tt">{head}</span>'
    if calc:
        inner += calc
    if beh:
        inner += f'<span class="tb">{beh}</span>'
    if watch:
        inner += f'<span class="tw">Watch: {watch}</span>'
    return f'<span class="mi-tip">{inner}</span>'


def hoverable(visible_html, tip_html, pos=""):
    """Wrap a value so hovering it shows its info card. pos: '' | 'mi-right' | 'mi-up'."""
    if not tip_html:
        return visible_html
    cls = ("mi-wrap " + pos).strip()
    return f'<span class="{cls}">{visible_html}{tip_html}</span>'


def _explain_block(label, items, dd, closed_idx):
    """A popover listing, for each (title, column, kind): the calculation + live behavior."""
    try:
        pop = st.popover(label)
    except Exception:
        pop = st.expander(label)
    with pop:
        for title, col, kind in items:
            calc, watch = METRIC_INFO.get(col, ("", ""))
            st.markdown(f"**{title}**")
            if calc:
                st.markdown(calc)
            beh = _behavior_html(dd, col, kind, closed_idx)
            if beh:
                st.markdown('<div class="nm-insight">' + beh + "</div>", unsafe_allow_html=True)
            if watch:
                st.caption("Watch: " + watch)
            st.divider()


def _chart_with_select(fig, key, config):
    """plotly_chart with click-selection when the Streamlit version supports it."""
    try:
        return st.plotly_chart(fig, use_container_width=True, key=key,
                               on_select="rerun", config=config)
    except TypeError:
        st.plotly_chart(fig, use_container_width=True, config=config)
        return None


@st.cache_data(ttl=86400, show_spinner=False)
def _gulf_geojson():
    """Country boundary polygons for the 5 markets (world.geo.json, ISO3 ids). Cached a day."""
    try:
        import requests
        r = requests.get(
            "https://raw.githubusercontent.com/johan/world.geo.json/master/countries.geo.json",
            timeout=10)
        gj = r.json()
        keep = set(COUNTRY_ISO.values())
        gj["features"] = [f for f in gj["features"] if f.get("id") in keep]
        return gj if gj["features"] else None
    except Exception:
        return None


def gulf_map(cts, vals, kind, title, height=470):
    """The market-intelligence-style map: country AREAS shaded by the metric, drawn over the
    Esri basemap (pan / scroll-zoom), value labels, no pins. Falls back to a plain filled map
    if the boundary file or tile service is unreachable."""
    go = _go()
    gj = _gulf_geojson()
    hover = [f"{c}: {fmt(v, kind)}" for c, v in zip(cts, vals)]
    labels = [f"<b>{c}</b><br>{fmt(v, kind)}" for c, v in zip(cts, vals)]
    lats = [COUNTRY_CENTROID[c][0] for c in cts]
    lons = [COUNTRY_CENTROID[c][1] for c in cts]
    scale = [[0.0, "#E5E9E0"], [0.55, "#5F8575"], [1.0, "#21362B"]]
    if gj is not None:
        try:
            _token = str(st.secrets.get("ARCGIS_API_KEY", "")).strip()
        except Exception:
            _token = ""
        _token = _token or ARCGIS_DEFAULT_KEY
        tile_url = ("https://static-map-tiles-api.arcgis.com/arcgis/rest/services/"
                    "static-basemap-tiles-service/v1/arcgis/navigation/static/tile/"
                    "{z}/{y}/{x}?token=" + _token)
        map_cfg = dict(style="white-bg", center=dict(lat=26.0, lon=49.6), zoom=4.3,
                       layers=[dict(below="traces", sourcetype="raster", source=[tile_url])])
        common = dict(geojson=gj, featureidkey="id",
                      locations=[COUNTRY_ISO[c] for c in cts],
                      z=[float(v) for v in vals], colorscale=scale, showscale=False,
                      marker_opacity=0.75, marker_line_color="white", marker_line_width=1.6,
                      hovertext=hover, hoverinfo="text", customdata=cts)
        try:
            fig = go.Figure(go.Choroplethmap(**common))
            fig.add_trace(go.Scattermap(lat=lats, lon=lons, mode="text", text=labels,
                                        textfont=dict(size=12, color="#21362B"),
                                        hoverinfo="skip", showlegend=False, customdata=cts))
            fig.update_layout(map=map_cfg)
        except (AttributeError, ValueError):
            fig = go.Figure(go.Choroplethmapbox(**common))
            fig.add_trace(go.Scattermapbox(lat=lats, lon=lons, mode="text", text=labels,
                                           textfont=dict(size=12, color="#21362B"),
                                           hoverinfo="skip", showlegend=False, customdata=cts))
            fig.update_layout(mapbox=map_cfg)
        fig.add_annotation(text="Powered by Esri", x=1, y=0, xref="paper", yref="paper",
                           showarrow=False, xanchor="right", yanchor="bottom",
                           font=dict(size=9, color="#7C776A"))
    else:
        fig = go.Figure(go.Choropleth(
            locations=[COUNTRY_ISO[c] for c in cts], z=[float(v) for v in vals],
            colorscale=scale, marker_line_color="white", marker_line_width=1.4,
            showscale=False, hovertext=hover, hoverinfo="text", customdata=cts))
        fig.add_trace(go.Scattergeo(lat=lats, lon=lons, mode="text", text=labels,
                                    textfont=dict(size=12, color="#21362B"),
                                    hoverinfo="skip", showlegend=False, customdata=cts))
        fig.update_geos(fitbounds="locations", visible=False, bgcolor="white",
                        showcountries=True, countrycolor="#DAD6C8",
                        showland=True, landcolor="#F8F7F2")
    fig.update_layout(
        title=dict(text=title, font=dict(size=14, color="#21362B", family="Arial Black, Arial")),
        height=height, margin=dict(l=4, r=4, t=44, b=4), paper_bgcolor="white",
        showlegend=False)
    return fig


def _clicked_country(ev):
    """Extract the clicked country from a plotly selection event (customdata carries it)."""
    try:
        pts = ev.selection.points
    except Exception:
        try:
            pts = ev["selection"]["points"]
        except Exception:
            return None
    if not pts:
        return None
    cd = pts[0].get("customdata")
    if isinstance(cd, (list, tuple)):
        cd = cd[0] if cd else None
    return str(cd) if cd else None


def donut(value, color, delta=None, height=168):
    """KPI ring with the percentage in the center (reference style)."""
    go = _go()
    v = float(value or 0)
    v01 = max(0.0, min(1.0, v))
    fig = go.Figure(go.Pie(values=[v01, 1 - v01], hole=0.74, sort=False, direction="clockwise",
                           marker=dict(colors=[color, LIGHT]), textinfo="none", hoverinfo="skip"))
    fig.add_annotation(text=f"<b>{v * 100:.0f}%</b>", showarrow=False,
                       font=dict(size=22, color="#21362B", family="Arial Black, Arial"))
    fig.update_layout(height=height, margin=dict(l=6, r=6, t=6, b=0),
                      showlegend=False, paper_bgcolor="white")
    return fig


def spark(series, color=TEAL, height=52):
    go = _go()
    fig = go.Figure(go.Scatter(y=list(series), mode="lines",
                               line=dict(color=color, width=2),
                               fill="tozeroy", fillcolor="rgba(33,54,43,0.10)"))
    fig.update_layout(height=height, margin=dict(l=0, r=0, t=2, b=0),
                      plot_bgcolor="white", paper_bgcolor="white", showlegend=False)
    fig.update_xaxes(visible=False)
    fig.update_yaxes(visible=False)
    return fig


def kpi_card(col, label, value_str, delta_val, delta_str, series, up_is_good=True,
             delta_label="vs prior", explain=None):
    with col:
        with st.container(border=True):
            st.markdown(f'<p class="kpi-l">{label}</p>', unsafe_allow_html=True)
            _tip = ""
            if explain and (explain[0] or explain[2]):
                _tin = f'<span class="tt">{label}</span>'
                if explain[0]:
                    _tin += explain[0]
                if explain[2]:
                    _tin += f'<span class="tb">{explain[2]}</span>'
                if explain[1]:
                    _tin += f'<span class="tw">Watch: {explain[1]}</span>'
                _tip = f'<span class="mi-tip">{_tin}</span>'
            st.markdown(f'<p class="kpi-v">{hoverable(value_str, _tip)}</p>',
                        unsafe_allow_html=True)
            if delta_val is None:
                st.markdown('<span class="kpi-d-na">no prior month</span>', unsafe_allow_html=True)
            else:
                good = (delta_val >= 0) if up_is_good else (delta_val <= 0)
                cls = "kpi-d-up" if good else "kpi-d-dn"
                arrow = "&#9650;" if delta_val >= 0 else "&#9660;"
                st.markdown(f'<span class="{cls}">{arrow} {delta_str} {delta_label}</span>',
                            unsafe_allow_html=True)
            if series is not None and len(series) > 1:
                st.plotly_chart(spark(series), use_container_width=True,
                                config={"displayModeBar": False, "staticPlot": True})
            if explain and (explain[0] or explain[2]):
                try:
                    with st.popover("calc & trend", use_container_width=True):
                        if explain[0]:
                            st.markdown(explain[0])
                        if explain[2]:
                            st.markdown('<div class="nm-insight">' + explain[2] + "</div>",
                                        unsafe_allow_html=True)
                        if explain[1]:
                            st.caption("Watch: " + explain[1])
                except Exception:
                    pass


def stacked_pct(d, x, buckets, title, closed_idx):
    """Colorful 100% stacked bars (reference style). buckets = [(label, col), ...]."""
    go = _go()
    fig = go.Figure()
    for (lbl, cn), color in zip(buckets, SERIES_COLORS):
        if cn in d.columns:
            fig.add_trace(go.Bar(x=x, y=d[cn], name=lbl, marker_color=color))
    fig.update_layout(barmode="stack")
    pct_axis(_base_layout(fig, title))
    return fig


# ---------------------------------------------------------------- All Hands scorecards
# Rows for the two Google Sheet "All Hands" scorecards. Tuples: (label, bridge column, format kind, up_is_good).
# label=None -> spacer row (blank line for visual grouping). up_is_good=True colors positive MoM green.
CK_SCORECARD_METRICS = [
    ("Kitchens sold",     "cws",                          "num",  True),
    ("Approved Deals",    "approved_deals",               "num",  True),
    ("RRA %",             "rra",                          "pct",  True),
    (None, None, None, None),
    ("Kitchens churned",  "churns_excl_transfers",        "num",  False),
    ("RRL %",             "rrl",                          "pct",  False),
    (None, None, None, None),
    ("Net Adds",          "net_adds",                     "num",  True),
    ("NRRA %",            "nrra",                         "pct",  True),
    (None, None, None, None),
    ("All AE Prod.",      "sales_team_cw_productivity",   "num1", True),
    (None, None, None, None),
    ("Owned sites",       "all_facilities",               "num",  True),
    ("Live kitchens",     "total_kitchens",               "num",  True),
    (None, None, None, None),
    ("Live Sold",         "live_sold_rate",               "pct",  True),
    ("Occupancy",         "occupancy",                    "pct",  True),
    (None, None, None, None),
    ("Occupied Kx",       "occupied_kitchens",            "num",  True),
]

CR_SCORECARD_METRICS = [
    ("Cloud Retail CWs",       "cr_cws",       "num", True),
    ("CR RRA $",               "cr_rra_usd",   "usd", True),
    (None, None, None, None),
    ("Cloud Retail Churns",    "cr_churns",    "num", False),
    ("CR RRL $",               "cr_rrl_usd",   "usd", False),
    (None, None, None, None),
    ("Cloud Retail Net Adds",  "cr_net_adds",  "num", True),  # calculated: cr_cws - cr_churns
    ("CR NRRA $",              "cr_nrra_usd",  "usd", True),
]

SCORECARD_REGIONS = ["Middle East", "UAE", "Saudi Arabia", "Kuwait"]
SCORECARD_REGION_DISPLAY = {"Saudi Arabia": "KSA"}


def _scorecard_cell_values(df, metrics_spec, win):
    """Prep the (values, mom, state) grid used by both the HTML render and the PPTX export.
    Returns (grid, df_maybe_augmented) where grid is a list of {label, kind, regions=[...]}
    or {spacer: True} entries."""
    if any(m[1] == "cr_net_adds" for m in metrics_spec if m[1]) \
            and "cr_net_adds" not in df.columns \
            and {"cr_cws", "cr_churns"} <= set(df.columns):
        df = df.copy()
        df["cr_net_adds"] = df["cr_cws"] - df["cr_churns"]
    grid = []
    for i, (label, col, kind, up_good) in enumerate(metrics_spec):
        if label is None:
            grid.append({"spacer": True})
            continue
        # A row "closes a group" when the very next entry is a spacer OR it is the last row.
        # Used to draw the horizontal underline between metric groups, matching the sheet.
        next_i = i + 1
        is_group_end = (
            next_i >= len(metrics_spec)
            or metrics_spec[next_i][0] is None
        )
        row = {"label": label, "kind": kind, "regions": [], "group_end": is_group_end}
        for region in SCORECARD_REGIONS:
            r = df[(df["country"] == region) & (df["month_end"].isin(win))]
            vals = []
            for m in win:
                rr = r[r["month_end"] == m]
                v = rr[col].iloc[0] if (col in r.columns and not rr.empty) else None
                try:
                    v = None if v is None or (isinstance(v, float) and math.isnan(v)) else float(v)
                except Exception:
                    v = None
                vals.append(v)
            mom = (vals[-1] - vals[-2]) if (vals[-1] is not None and vals[-2] is not None) else None
            if mom is None:
                state = "na"
            elif abs(mom) < 1e-12:
                state = "flat"
            else:
                good = (mom > 0) if up_good else (mom < 0)
                state = "up" if good else "dn"
            row["regions"].append({"vals": vals, "mom": mom, "state": state})
        grid.append(row)
    return grid, df


def _pptx_style_cell(cell, *, bold=False, size=10, align="left", color=(0x21, 0x36, 0x2B), fill=None):
    """Apply font + alignment + fill to a python-pptx table cell. All three color styles
    (bold, size, color, alignment) must be applied AFTER the cell.text assignment because
    setting .text discards any previous run styling."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    tf = cell.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    align_map = {"left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT, "center": PP_ALIGN.CENTER}
    for p in tf.paragraphs:
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        for r in p.runs:
            r.font.bold = bool(bold)
            r.font.size = Pt(size)
            r.font.color.rgb = RGBColor(*color)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*fill)


def _pptx_clear_table_style(tbl):
    """Strip PowerPoint's default table style off a python-pptx table so the theme's
    Medium Style (blue banded rows, blue header, grid borders) does not paint over
    our explicit cell fills and borders. Sets the tableStyleId to the 'no style, no
    grid' preset and turns off first-row / banded-row / etc. flags on tblPr."""
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tblPr = tbl._tbl.tblPr
    for old in tblPr.findall(f"{{{ns_a}}}tableStyleId"):
        tblPr.remove(old)
    style_id = etree.SubElement(tblPr, f"{{{ns_a}}}tableStyleId")
    style_id.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid
    for attr in ("firstRow", "firstCol", "lastRow", "lastCol", "bandRow", "bandCol"):
        tblPr.attrib.pop(attr, None)


def _pptx_remove_cell_borders(cell):
    """Explicitly turn off left/right/top/bottom borders on a table cell so the
    inherited table style doesn't draw grid lines. Bottom-border helper below is
    then free to add a single visible line where we want the group-end rule."""
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tc = cell._tc
    tcPr = tc.find(f"{{{ns_a}}}tcPr")
    if tcPr is None:
        tcPr = etree.SubElement(tc, f"{{{ns_a}}}tcPr")
    for side in ("lnL", "lnR", "lnT", "lnB"):
        for old in tcPr.findall(f"{{{ns_a}}}{side}"):
            tcPr.remove(old)
        ln = etree.SubElement(tcPr, f"{{{ns_a}}}{side}")
        ln.set("w", "6350")
        etree.SubElement(ln, f"{{{ns_a}}}noFill")


def _pptx_cell_fill(cell, rgb):
    """Solid fill on a python-pptx table cell."""
    from pptx.dml.color import RGBColor
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(*rgb)


def _pptx_cell_bottom_border(cell, *, color_hex="B7B0A0", width_emu=6350):
    """Add a bottom border to a python-pptx table cell. python-pptx has no first-class
    cell-border API, so we build the OOXML fragment (a:lnB) directly on the cell's tcPr."""
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    tc = cell._tc
    tcPr = tc.find(f"{{{ns_a}}}tcPr")
    if tcPr is None:
        tcPr = etree.SubElement(tc, f"{{{ns_a}}}tcPr")
    # Remove any prior bottom border so we do not stack styles.
    for old in tcPr.findall(f"{{{ns_a}}}lnB"):
        tcPr.remove(old)
    lnB = etree.SubElement(tcPr, f"{{{ns_a}}}lnB")
    lnB.set("w", str(width_emu))
    lnB.set("cap", "flat")
    lnB.set("cmpd", "sng")
    lnB.set("algn", "ctr")
    solid = etree.SubElement(lnB, f"{{{ns_a}}}solidFill")
    srgb = etree.SubElement(solid, f"{{{ns_a}}}srgbClr")
    srgb.set("val", color_hex)


def _brand_slide(slide, prs, *, slide_number=None):
    """Paint the NAMAA deck chrome on a blank slide: cream background, NAMAA logo
    top-left, terracotta motif bottom-right, confidential footer bottom-left,
    slide number bottom-right. Matches the reference deck template."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # Background: NAMAA cream. python-pptx exposes slide.background.fill,
    # so no XML mangling is needed for this.
    try:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xEE, 0xED, 0xE5)
    except Exception:
        # Some template variants block programmatic background changes; fall back to
        # a full-slide cream rectangle behind everything else.
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_w, slide_h)
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xEE, 0xED, 0xE5)
        rect.line.fill.background()
        # Send behind everything else (best-effort).
        try:
            spTree = rect._element.getparent()
            spTree.remove(rect._element)
            spTree.insert(2, rect._element)
        except Exception:
            pass

    # NAMAA logo top-left. Height is roughly the header band; width auto-scales.
    if os.path.exists(LOGO_PATH):
        try:
            slide.shapes.add_picture(
                LOGO_PATH,
                left=Inches(0.45),
                top=Inches(0.35),
                height=Inches(0.55),
            )
        except Exception:
            pass

    # Terracotta motif bottom-right - subtle brand accent from the reference deck.
    if os.path.exists(MOTIF_PATH):
        try:
            motif_w = Inches(4.6)
            slide.shapes.add_picture(
                MOTIF_PATH,
                left=slide_w - motif_w + Inches(0.25),
                top=slide_h - Inches(3.4),
                width=motif_w,
            )
        except Exception:
            pass

    # Bottom-left: confidential footer, small dark green text.
    footer_h = Inches(0.35)
    tb = slide.shapes.add_textbox(Inches(0.45), slide_h - footer_h - Inches(0.15),
                                  Inches(8), footer_h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "INFORMATION IN THIS DOCUMENT IS PRIVATE AND CONFIDENTIAL"
    r.font.size = Pt(9)
    r.font.bold = False
    r.font.color.rgb = RGBColor(0x21, 0x36, 0x2B)

    # Bottom-right: slide number (or blank if unset).
    if slide_number is not None:
        tb = slide.shapes.add_textbox(slide_w - Inches(0.85), slide_h - footer_h - Inches(0.15),
                                      Inches(0.5), footer_h)
        tf = tb.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = str(slide_number)
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(0x21, 0x36, 0x2B)


def _pptx_table_no_grid(table):
    """Kill PowerPoint's default table chrome (banding + the vertical gridlines):
    swap the table style to the built-in 'No Style, No Grid' and disable banding."""
    from pptx.oxml.ns import qn
    tbl_el = table._tbl
    tblPr = tbl_el.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("firstRow", "0")
    for el in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(el)
    style_el = tblPr.makeelement(qn("a:tableStyleId"), {})
    style_el.text = "{5940675A-B579-460E-94D1-54222C63F5DA}"  # No Style, No Grid
    tblPr.append(style_el)


def _pptx_cell_no_vlines(cell):
    """Belt-and-braces: explicit no-fill left/right borders so no theme can draw verticals."""
    from pptx.oxml.ns import qn
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnR", "a:lnL"):   # insert lnL last so final order is lnL, lnR (schema order)
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
        ln = tcPr.makeelement(qn(tag), {"w": "0"})
        ln.append(tcPr.makeelement(qn("a:noFill"), {}))
        tcPr.insert(0, ln)


def _build_scorecard_pptx(title, sel_label, win_hdrs, grid) -> bytes:
    """Native-PowerPoint slide with a wide table matching the on-screen scorecard.
    Editable in PowerPoint (not an image). Widescreen 16:9 layout with NAMAA brand chrome."""
    from io import BytesIO
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Paint NAMAA background + logo + motif + footer FIRST so the table renders on top.
    _brand_slide(slide, prs, slide_number=11)

    # Slide title (dark green, right of the logo).
    tb = slide.shapes.add_textbox(Inches(2.0), Inches(0.4), Inches(11.0), Inches(0.55))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{title}  —  {sel_label}"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x21, 0x36, 0x2B)

    n_regions = len(SCORECARD_REGIONS)
    n_cols = 1 + n_regions * 4
    body_rows = [g for g in grid if not g.get("spacer")]
    n_rows = 2 + len(body_rows)
    left, top = Inches(0.35), Inches(1.15)
    width = Inches(12.7)
    height = Inches(0.55) + Inches(0.32) * len(body_rows)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    # Strip PowerPoint's default table style so the theme doesn't paint blue banded
    # rows and grid borders over our explicit cell fills.
    try:
        _pptx_clear_table_style(tbl)
    except Exception:
        pass
    try:
        _pptx_table_no_grid(tbl)  # legacy no-op helper — keep alongside for older PPT viewers
    except Exception:
        pass

    # Column widths: widen the MoM column so "MoM Δ" and values like "-6.9%" fit on one line.
    label_w_in = 1.4
    month_w_in = 0.58
    mom_w_in = 1.00
    tbl.columns[0].width = Inches(label_w_in)
    ci = 1
    for _ in SCORECARD_REGIONS:
        for _ in range(3):
            tbl.columns[ci].width = Inches(month_w_in)
            ci += 1
        tbl.columns[ci].width = Inches(mom_w_in)
        ci += 1

    cream_bg = (0xEE, 0xED, 0xE5)
    body_bg = (0xF5, 0xF1, 0xE7)
    banner_bg = (0x21, 0x36, 0x2B)

    # Row 0: region banner (merged across the 4 cols per region), transparent cream fill.
    tbl.cell(0, 0).text = ""
    _pptx_style_cell(tbl.cell(0, 0), size=10)
    _pptx_cell_fill(tbl.cell(0, 0), cream_bg)
    _pptx_remove_cell_borders(tbl.cell(0, 0))
    for r_i, region in enumerate(SCORECARD_REGIONS):
        c0 = 1 + r_i * 4
        c1 = c0 + 3
        tbl.cell(0, c0).merge(tbl.cell(0, c1))
        tbl.cell(0, c0).text = SCORECARD_REGION_DISPLAY.get(region, region)
        _pptx_style_cell(tbl.cell(0, c0), bold=True, size=13, align="center",
                         color=(0x21, 0x36, 0x2B))
        _pptx_cell_fill(tbl.cell(0, c0), cream_bg)
        _pptx_remove_cell_borders(tbl.cell(0, c0))

    # Row 1: month subheaders + MoM Δ, dark green banner, white text.
    tbl.cell(1, 0).text = ""
    _pptx_style_cell(tbl.cell(1, 0), size=9)
    _pptx_cell_fill(tbl.cell(1, 0), cream_bg)
    _pptx_remove_cell_borders(tbl.cell(1, 0))
    ci = 1
    for _ in SCORECARD_REGIONS:
        for h in win_hdrs:
            tbl.cell(1, ci).text = h
            _pptx_style_cell(tbl.cell(1, ci), bold=True, size=10, align="center",
                             fill=banner_bg, color=(0xFF, 0xFF, 0xFF))
            _pptx_remove_cell_borders(tbl.cell(1, ci))
            ci += 1
        tbl.cell(1, ci).text = "MoM Δ"
        _pptx_style_cell(tbl.cell(1, ci), bold=True, size=10, align="center",
                         fill=banner_bg, color=(0xFF, 0xFF, 0xFF))
        _pptx_remove_cell_borders(tbl.cell(1, ci))
        ci += 1

    up_rgb = (0x16, 0xA3, 0x4A)
    dn_rgb = (0xB4, 0x47, 0x2E)
    flat_rgb = (0x64, 0x74, 0x8B)
    row_idx = 2
    for g in grid:
        if g.get("spacer"):
            continue
        # Label cell.
        tbl.cell(row_idx, 0).text = g["label"]
        _pptx_style_cell(tbl.cell(row_idx, 0), bold=True, size=10, align="left",
                         color=(0x21, 0x36, 0x2B))
        _pptx_cell_fill(tbl.cell(row_idx, 0), cream_bg)
        _pptx_remove_cell_borders(tbl.cell(row_idx, 0))
        ci = 1
        for reg in g["regions"]:
            # Value cells.
            for v in reg["vals"]:
                tbl.cell(row_idx, ci).text = fmt(v, g["kind"]) if v is not None else "—"
                _pptx_style_cell(tbl.cell(row_idx, ci), size=10, align="right",
                                 color=(0x21, 0x36, 0x2B))
                _pptx_cell_fill(tbl.cell(row_idx, ci), body_bg)
                _pptx_remove_cell_borders(tbl.cell(row_idx, ci))
                if g.get("group_end"):
                    try:
                        _pptx_cell_bottom_border(tbl.cell(row_idx, ci))
                    except Exception:
                        pass
                ci += 1
            # MoM cell.
            state = reg["state"]
            mom_txt = "—" if reg["mom"] is None else fmt(reg["mom"], g["kind"])
            tbl.cell(row_idx, ci).text = mom_txt
            mom_color = up_rgb if state == "up" else dn_rgb if state == "dn" else flat_rgb
            _pptx_style_cell(tbl.cell(row_idx, ci), bold=True, size=10, align="right",
                             color=mom_color)
            _pptx_cell_fill(tbl.cell(row_idx, ci), body_bg)
            _pptx_remove_cell_borders(tbl.cell(row_idx, ci))
            if g.get("group_end"):
                try:
                    _pptx_cell_bottom_border(tbl.cell(row_idx, ci))
                except Exception:
                    pass
            ci += 1
        row_idx += 1

    # Note: cell borders are already stripped inline via _pptx_remove_cell_borders() during
    # the row loop, and the group-end bottom borders are applied after that on specific rows.
    # An additional post-loop sweep with the legacy _pptx_cell_no_vlines helper would risk
    # wiping the group-end bottom border on newer python-pptx versions, so we skip it.

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _render_all_hands_scorecard(df, all_months, metrics_spec, key_prefix, *, pptx_title=None):
    """Reproduce the Google Sheet All Hands scorecard layout: 4 region blocks side-by-side,
    each with the last 3 months (as-of picker) + a MoM Δ column, one row per metric.
    Negative MoM in red / positive in green (unless up_is_good=False, then flipped)."""
    if df is None or df.empty or not all_months:
        st.info("No panel data available.")
        return

    # As-of month picker (newest first) + PPTX download side-by-side.
    month_labels = [pd.Timestamp(m).strftime("%b %Y") for m in all_months]
    labels_desc = list(reversed(month_labels))
    picker_col, dl_col = st.columns([3, 1], vertical_alignment="bottom")
    with picker_col:
        sel_label = st.selectbox(
            "As-of month",
            options=labels_desc,
            index=0,
            key=f"{key_prefix}_scorecard_asof",
        )
    sel_idx = month_labels.index(sel_label)
    if sel_idx < 2:
        st.warning("Need at least 3 months of history to build the scorecard. "
                   "Pick a more recent as-of month or wait for more data.")
        return
    win = all_months[sel_idx - 2: sel_idx + 1]  # oldest, mid, newest (as-of)
    win_hdrs = [pd.Timestamp(m).strftime("%b %y") for m in win]

    # Build the grid once - drives both the HTML render and the PPTX export.
    grid, df = _scorecard_cell_values(df, metrics_spec, win)

    # Download-as-PowerPoint button.
    with dl_col:
        title_for_pptx = pptx_title or "All Hands"
        try:
            pptx_bytes = _build_scorecard_pptx(title_for_pptx, sel_label, win_hdrs, grid)
        except Exception as e:
            pptx_bytes = None
            st.caption(f"PPTX export unavailable: {type(e).__name__}")
        if pptx_bytes:
            fname_safe = title_for_pptx.replace(" ", "_")
            month_slug = sel_label.replace(" ", "_")
            st.download_button(
                "Download as slide",
                data=pptx_bytes,
                file_name=f"{fname_safe}_{month_slug}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key=f"{key_prefix}_scorecard_pptx",
                use_container_width=True,
            )

    # Build the HTML table for the on-screen render.
    n_regions = len(SCORECARD_REGIONS)
    total_cols = 1 + n_regions * 4
    html = [
        "<div class='scorecard-wrap'>",
        f"<div class='scorecard-title'>{sel_label}</div>",
        "<table class='scorecard'>",
        "<colgroup>",
        "<col class='c-label'/>",
    ]
    for _ in SCORECARD_REGIONS:
        html += ["<col/>", "<col/>", "<col/>", "<col class='c-delta'/>"]
    html.append("</colgroup>")

    html.append("<thead>")
    html.append("<tr>")
    html.append("<th class='sc-empty'></th>")
    for r in SCORECARD_REGIONS:
        html.append(f"<th class='sc-region' colspan='4'>{SCORECARD_REGION_DISPLAY.get(r, r)}</th>")
    html.append("</tr>")
    html.append("<tr>")
    html.append("<th class='sc-empty'></th>")
    for _ in SCORECARD_REGIONS:
        for h in win_hdrs:
            html.append(f"<th class='sc-month'>{h}</th>")
        html.append("<th class='sc-month'>MoM Δ</th>")
    html.append("</tr>")
    html.append("</thead><tbody>")

    state_to_cls = {"up": "sc-delta-up", "dn": "sc-delta-dn",
                    "flat": "sc-delta-flat", "na": "sc-delta-na"}
    for g in grid:
        if g.get("spacer"):
            html.append(f"<tr class='sc-spacer'><td colspan='{total_cols}'>&nbsp;</td></tr>")
            continue
        row_cls = "sc-group-end" if g.get("group_end") else ""
        html.append(f"<tr class='{row_cls}'>")
        html.append(f"<td class='sc-metric'>{g['label']}</td>")
        for reg in g["regions"]:
            for v in reg["vals"]:
                html.append(f"<td class='sc-val'>{fmt(v, g['kind']) if v is not None else '—'}</td>")
            mom_txt = "—" if reg["mom"] is None else fmt(reg["mom"], g["kind"])
            html.append(f"<td class='sc-delta {state_to_cls[reg['state']]}'>{mom_txt}</td>")
        html.append("</tr>")

    html.append("</tbody></table></div>")

    st.markdown(
        """
        <style>
        .scorecard-wrap { overflow-x: auto; padding: 6px 0 18px 0; }
        .scorecard-title { font-size: 1.1rem; font-weight: 800; color: #1F3B57; margin: 6px 0 10px 2px; }
        table.scorecard { border-collapse: separate; border-spacing: 0; width: 100%;
                          font-family: Arial, sans-serif; font-size: 0.85rem; background: #F5F1E7; }
        table.scorecard th, table.scorecard td { padding: 6px 10px; text-align: right; white-space: nowrap; }
        table.scorecard col.c-label { width: 190px; }
        table.scorecard col.c-delta { width: 82px; }
        table.scorecard th.sc-empty { background: transparent; }
        table.scorecard th.sc-region { background: transparent; color: #1F3B57; font-weight: 800;
                                       text-align: center; font-size: 1rem; padding-bottom: 4px; }
        table.scorecard th.sc-month { background: #21362B; color: #FFFFFF; font-weight: 700;
                                      text-align: center; font-size: 0.78rem; letter-spacing: 0.02em; }
        table.scorecard td.sc-metric { text-align: left; font-weight: 700; color: #21362B; }
        table.scorecard td.sc-val { color: #21362B; font-weight: 600; }
        table.scorecard td.sc-delta { font-weight: 800; }
        table.scorecard .sc-delta-up { color: #16A34A; }
        table.scorecard .sc-delta-dn { color: #E74C3C; }
        table.scorecard .sc-delta-flat { color: #64748B; }
        table.scorecard .sc-delta-na { color: #94A3B8; }
        /* Group-end rows get a thin bottom border under the value + MoM cells (like the sheet).
           The label cell is intentionally left borderless to keep the visual grouping clean. */
        table.scorecard tr.sc-group-end td.sc-val,
        table.scorecard tr.sc-group-end td.sc-delta {
            border-bottom: 1px solid #B7B0A0;
            padding-bottom: 8px;
        }
        table.scorecard tr.sc-spacer td { padding: 0; height: 10px; background: transparent; }
        </style>
        """ + "".join(html),
        unsafe_allow_html=True,
    )


# ---------------------------------------------------------------- Panel Overview (existing dashboard)
# ---------------------------------------------------------------- Executive view
# The CEO/GM mental model, taken from the All Hands slides: all four markets side by side,
# the slide's exact metric groups (Growth -> Churn -> Net -> Productivity -> Footprint ->
# Utilization), MoM direction on everything, %-of-base revenue metrics rather than raw $.
EXEC_GROUPS = [
    ("Growth", [("Kitchens sold", "cws", "num", True),
                ("Approved Deals", "approved_deals", "num", True),
                ("RRA %", "rra", "pct", True)]),
    ("Churn", [("Kitchens churned", "churns_excl_transfers", "num", False),
               ("RRL %", "rrl", "pct", False)]),
    ("Net", [("Net Adds", "net_adds", "num", True),
             ("NRRA %", "nrra", "pct", True)]),
    ("Productivity", [("AE Productivity (CWs/AE)", "sales_team_cw_productivity", "num1", True)]),
    ("Footprint", [("Owned sites", "all_facilities", "num", True),
                   ("Live kitchens", "total_kitchens", "num", True)]),
    ("Utilization", [("Live Sold", "live_sold_rate", "pct", True),
                     ("Occupancy", "occupancy", "pct", True),
                     ("Occupied kitchens", "occupied_kitchens", "num", True)]),
]
EXEC_REGIONS = ["Middle East", "UAE", "Saudi Arabia", "Kuwait", "Bahrain", "Qatar"]
EXEC_RDISP = {"Saudi Arabia": "KSA", "Middle East": "ME", "Bahrain": "BH", "Qatar": "QAT"}
EXEC_REG_COLORS = {"Middle East": NAVY, "UAE": ORANGE, "Saudi Arabia": TEAL,
                   "Kuwait": YELLOW, "Bahrain": RED, "Qatar": SLATE}


def _render_executive_view(df, all_months, cur_month_start):
    months_closed = [m for m in all_months if pd.Timestamp(m) < cur_month_start]
    if len(months_closed) < 2:
        st.info("Not enough closed months for the executive view.")
        return
    # ---- executive controls: everything below follows these ----
    _lbls = [pd.Timestamp(m).strftime("%b %Y") for m in months_closed]
    ec1, ec2 = st.columns([1.6, 3.4], vertical_alignment="bottom")
    with ec1:
        _sel_asof = st.selectbox("As of month", list(reversed(_lbls)), index=0, key="exec_asof")
    _ai = _lbls.index(_sel_asof)
    asof = months_closed[_ai]
    prev = months_closed[_ai - 1] if _ai >= 1 else None
    asof_lbl = pd.Timestamp(asof).strftime("%b %Y")
    with ec2:
        try:
            _win_sel = st.segmented_control("Momentum window", ["3M", "6M", "12M"],
                                            default="6M", key="exec_win") or "6M"
        except Exception:
            _win_sel = st.radio("Momentum window", ["3M", "6M", "12M"], index=1,
                                horizontal=True, key="exec_win")
    _wn = int(_win_sel[:-1])
    win_months = months_closed[max(0, _ai - _wn + 1): _ai + 1]

    def val(region, col, m):
        r = df[(df["country"] == region) & (df["month_end"] == m)]
        if r.empty or col not in r.columns:
            return None
        try:
            v = float(r[col].iloc[0])
            return None if math.isnan(v) else v
        except Exception:
            return None

    def delta(region, col):
        a, b = val(region, col, asof), val(region, col, prev)
        return (a - b) if (a is not None and b is not None) else None

    _dd_cache = {}

    def _dd_ci(region):
        """Region series (sorted, labeled) + index of the selected as-of month."""
        if region not in _dd_cache:
            d = df[df["country"] == region].sort_values("month_end").reset_index(drop=True)
            d["month_label"] = d["month_end"].dt.strftime("%b %Y")
            idx = d.index[d["month_end"] == asof]
            _dd_cache[region] = (d, int(idx[0]) if len(idx) else max(0, len(d) - 1))
        return _dd_cache[region]

    def tip(region, label, col, kind):
        d, ci = _dd_ci(region)
        return _metric_tip(label, col, kind, d, ci, scope=EXEC_RDISP.get(region, region))

    def arrow_html(dv, kind, up_good):
        if dv is None:
            return '<span class="kpi-d-na">-</span>'
        good = (dv >= 0) if up_good else (dv <= 0)
        cls = "kpi-d-up" if good else "kpi-d-dn"
        ar = "&#9650;" if dv >= 0 else "&#9660;"
        dtxt = f"{dv * 100:+.1f} pp" if kind == "pct" else ("+" if dv >= 0 else "") + fmt(dv, kind)
        return f'<span class="{cls}">{ar} {dtxt}</span>'

    # ---- headline narrative ----
    _parts = [f"As of <b>{asof_lbl}</b> (last closed month)"]
    _na, _dna = val("Middle East", "net_adds", asof), delta("Middle East", "net_adds")
    if _na is not None:
        _s = f"ME net adds <b>{fmt(_na, 'num')}</b>"
        if _dna is not None:
            _s += f" ({'+' if _dna >= 0 else ''}{fmt(_dna, 'num')} MoM)"
        _parts.append(_s)
    _mkts = [(r, val(r, "net_adds", asof)) for r in EXEC_REGIONS[1:]]
    _mkts = [(r, v) for r, v in _mkts if v is not None]
    if _mkts:
        _bst = max(_mkts, key=lambda t: t[1])
        _wst = min(_mkts, key=lambda t: t[1])
        _parts.append(f"strongest market: <b>{EXEC_RDISP.get(_bst[0], _bst[0])}</b> ({fmt(_bst[1], 'num')} net adds)")
        _parts.append(f"weakest: <b>{EXEC_RDISP.get(_wst[0], _wst[0])}</b> ({fmt(_wst[1], 'num')})")
    insight(" &middot; ".join(_parts))

    # ---- AI brief + anomaly radar ----
    _anoms = _anomaly_scan(df, asof, months_closed)
    _bc1, _bc2 = st.columns([1.6, 4.4], vertical_alignment="center")
    with _bc1:
        _gen_brief = st.button(":sparkles: Generate AI brief", key="exec_ai_brief_btn",
                               use_container_width=True)
    with _bc2:
        if _anoms:
            _chips = []
            for a in _anoms:
                _acls = "hot" if abs(a["z"]) >= 3 else "warm"
                _dt = (f"{a['delta'] * 100:+.1f} pp" if a["kind"] == "pct"
                       else ("+" if a["delta"] >= 0 else "") + fmt(a["delta"], a["kind"]))
                _chips.append(f'<span class="anom {_acls}">'
                              f'{EXEC_RDISP.get(a["region"], a["region"])} &middot; '
                              f'{a["metric"]} {_dt} (z {a["z"]:+.1f})</span>')
            st.markdown('<div class="anom-row"><span class="anom-t">Anomaly radar</span>'
                        + "".join(_chips) + "</div>", unsafe_allow_html=True)

    def _offline_brief():
        na, dna = val("Middle East", "net_adds", asof), delta("Middle East", "net_adds")
        L = ["**Headline**",
             f"ME closed {asof_lbl} at **{fmt(na, 'num') or '-'} net adds**"
             + (f" ({'+' if dna >= 0 else ''}{fmt(dna, 'num')} MoM)" if dna is not None else "")
             + ".", "", "**What moved**"]
        for _lbl, _cn, _k in [("Kitchens sold", "cws", "num"),
                              ("Kitchens churned", "churns_excl_transfers", "num"),
                              ("NRRA", "nrra", "pct"), ("Occupancy", "occupancy", "pct")]:
            _v, _d = val("Middle East", _cn, asof), delta("Middle East", _cn)
            if _v is None:
                continue
            _s = f"- {_lbl}: **{fmt(_v, _k)}**"
            if _d is not None:
                _s += (f" ({_d * 100:+.1f} pp MoM)" if _k == "pct"
                       else f" ({'+' if _d >= 0 else ''}{fmt(_d, _k)} MoM)")
            L.append(_s)
        if _mkts:
            L.append(f"- Strongest market {EXEC_RDISP.get(_bst[0], _bst[0])}"
                     f" ({fmt(_bst[1], 'num')} net adds); weakest {EXEC_RDISP.get(_wst[0], _wst[0])}"
                     f" ({fmt(_wst[1], 'num')}).")
        if _anoms:
            L += ["", "**Watch-outs**"]
            L += [f"- {EXEC_RDISP.get(a['region'], a['region'])} {a['metric']}: unusually large "
                  f"move this month (z {a['z']:+.1f})." for a in _anoms[:3]]
        return "\n".join(L)

    _bk = "exec_brief_" + asof_lbl + "_" + _win_sel
    if _gen_brief:
        with st.spinner("Writing the brief..."):
            _btxt = None
            if _ai_key():
                try:
                    _bprompt = (
                        f"Write the executive brief for {asof_lbl} for the Middle East business. "
                        "Sections: **Headline** (one sentence), **What moved** (3-5 bullets with "
                        "numbers, naming markets), **Watch-outs** (2-3 bullets), **Do next** (2 "
                        "bullets). Maximum 180 words. No chart block."
                        + (" Known anomalies: " + "; ".join(
                            f"{a['region']} {a['metric']} z {a['z']:+.1f}" for a in _anoms)
                           if _anoms else ""))
                    _braw = "".join(_claude_stream(
                        [{"role": "user", "content": _bprompt}],
                        _ai_system(df, cur_month_start)))
                    _btxt, _ = _split_ai_answer(_braw)
                except Exception as _be:
                    st.caption("Claude call failed (" + str(_be)[:100] + ") - offline brief below.")
            if not _btxt:
                _btxt = _offline_brief()
            st.session_state[_bk] = _btxt
    if st.session_state.get(_bk):
        with st.container(border=True):
            st.markdown(st.session_state[_bk])
            st.caption("NAMAA Copilot brief - " + asof_lbl
                       + (" - Claude" if _ai_key() else " - offline engine"))

    # ---- market summary cards ----
    cards = st.columns(len(EXEC_REGIONS))
    for region, c in zip(EXEC_REGIONS, cards):
        with c:
            with st.container(border=True):
                st.markdown(f'<p class="kpi-l">{EXEC_RDISP.get(region, region)}</p>',
                            unsafe_allow_html=True)
                _v = val(region, "net_adds", asof)
                _pos = "mi-right" if region in EXEC_REGIONS[-2:] else ""
                st.markdown('<p class="kpi-v">'
                            + hoverable(fmt(_v, "num") or "-",
                                        tip(region, "Net Adds", "net_adds", "num"), _pos)
                            + ' <span style="font-size:0.8rem; '
                            'font-weight:600; color:#7C776A;">net adds</span></p>',
                            unsafe_allow_html=True)
                st.markdown(arrow_html(delta(region, "net_adds"), "num", True) + " MoM",
                            unsafe_allow_html=True)
                _rows = []
                for lbl, cn, kind in [("NRRA %", "nrra", "pct"), ("Live Sold", "live_sold_rate", "pct"),
                                      ("Occupancy", "occupancy", "pct")]:
                    _vv = val(region, cn, asof)
                    if _vv is not None:
                        _rows.append(hoverable(f'{lbl} <b>{fmt(_vv, kind)}</b>',
                                               tip(region, lbl, cn, kind), _pos)
                                     + f' {arrow_html(delta(region, cn), kind, True)}')
                if _rows:
                    st.markdown('<div style="font-size:0.8rem; color:#4A5548; margin-top:4px;">'
                                + "<br>".join(_rows) + "</div>", unsafe_allow_html=True)
                if st.button("Deep dive", key=f"exec_dd_{region}", use_container_width=True):
                    st.session_state["_pending_cty"] = region
                    st.session_state["_pending_view"] = "Market deep dive"
                    st.rerun()

    # ---- the slide grid: metric x market with MoM arrows ----
    sec("Scorecard", f"The All Hands metric set - every market at once, {asof_lbl} with MoM direction")
    _h = ['<div class="exec-wrap"><table class="exec-grid"><tr><th>Metric</th>']
    for region in EXEC_REGIONS:
        _h.append(f"<th>{EXEC_RDISP.get(region, region)}</th>")
    _h.append("</tr>")
    for gname, items in EXEC_GROUPS:
        _h.append(f'<tr class="exec-band"><td colspan="{len(EXEC_REGIONS) + 1}">{gname}</td></tr>')
        for lbl, cn, kind, up_good in items:
            if cn not in df.columns:
                continue
            _h.append(f"<tr><td>{hoverable(lbl, tip('Middle East', lbl, cn, kind))}</td>")
            for region in EXEC_REGIONS:
                _v = val(region, cn, asof)
                _cell = f'<span class="ev">{fmt(_v, kind) if _v is not None else "-"}</span>'
                _cell += f'<span class="ed">{arrow_html(delta(region, cn), kind, up_good)}</span>'
                _rpos = "mi-right" if region in EXEC_REGIONS[-3:] else ""
                _h.append(f"<td>{hoverable(_cell, tip(region, lbl, cn, kind), _rpos)}</td>")
            _h.append("</tr>")
    _h.append("</table></div>")
    st.markdown("".join(_h), unsafe_allow_html=True)

    # Per-metric calculation + live behavior (ME series) for everything on the scorecard.
    dd_me = df[df["country"] == "Middle East"].sort_values("month_end").reset_index(drop=True)
    if "month_label" not in dd_me.columns:
        dd_me["month_label"] = dd_me["month_end"].dt.strftime("%b %Y")
    _ci = max(0, len([m for m in dd_me["month_end"] if pd.Timestamp(m) < cur_month_start]) - 1)
    _explain_block("How each metric is calculated & behaving (ME)",
                   [(lbl, cn, kind) for _, items in EXEC_GROUPS
                    for (lbl, cn, kind, _) in items if cn in df.columns],
                   dd_me, _ci)

    # ---- metric drill-down: tap any scorecard metric to expand it across markets ----
    sec("Drill-down", "Tap a metric - trend by market, ranking, and how it is calculated")
    _flatm = [(lbl, cn, kind, ug) for _, items in EXEC_GROUPS
              for (lbl, cn, kind, ug) in items if cn in df.columns]
    _dl_lbls = [t[0] for t in _flatm]
    try:
        _dsel = st.pills("Metric", _dl_lbls, default="Net Adds" if "Net Adds" in _dl_lbls else _dl_lbls[0],
                         key="exec_drill") or _dl_lbls[0]
    except Exception:
        try:
            _dsel = st.segmented_control("Metric", _dl_lbls,
                                         default="Net Adds" if "Net Adds" in _dl_lbls else _dl_lbls[0],
                                         key="exec_drill") or _dl_lbls[0]
        except Exception:
            _dsel = st.selectbox("Metric", _dl_lbls, key="exec_drill")
    _dlbl, _dcol, _dkind, _dug = next((t for t in _flatm if t[0] == _dsel), _flatm[0])
    go = _go()
    dr1, dr2 = st.columns([3, 2])
    with dr1:
        fig = go.Figure()
        for region in EXEC_REGIONS:
            ys = [val(region, _dcol, m) for m in win_months]
            fig.add_trace(go.Scatter(
                x=[pd.Timestamp(m).strftime("%b %y") for m in win_months], y=ys,
                mode="lines+markers", name=EXEC_RDISP.get(region, region),
                line=dict(color=EXEC_REG_COLORS.get(region, SLATE), width=2.4),
                marker=dict(size=7, color="white",
                            line=dict(color=EXEC_REG_COLORS.get(region, SLATE), width=2)),
                customdata=[fmt(v, _dkind) for v in ys],
                hovertemplate="<b>%{fullData.name}</b>: %{customdata}<extra></extra>"))
        _base_layout(fig, f"{_dlbl} - trend by market", height=360)
        if _dkind in ("usd", "usdk"):
            money_axis(fig)
        elif _dkind == "pct":
            pct_axis(fig)
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with dr2:
        _rank = [(r, val(r, _dcol, asof)) for r in EXEC_REGIONS[1:]]
        _rank = sorted([(r, v) for r, v in _rank if v is not None], key=lambda t: t[1])
        if _rank:
            fig = go.Figure(go.Bar(
                x=[v for _, v in _rank], y=[EXEC_RDISP.get(r, r) for r, _ in _rank],
                orientation="h",
                marker_color=[EXEC_REG_COLORS.get(r, SLATE) for r, _ in _rank],
                text=[fmt(v, _dkind) for _, v in _rank], textposition="outside"))
            _base_layout(fig, f"{_dlbl} - {asof_lbl}", height=360)
            fig.update_xaxes(showticklabels=False)
            fig.update_layout(margin=dict(l=8, r=70, t=44, b=4))
            st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    _calc, _watch = METRIC_INFO.get(_dcol, ("", ""))
    _beh = _behavior_html(dd_me, _dcol, _dkind, _ci)
    if _calc or _beh:
        _txt = (_calc + (" &middot; " if _calc and _beh else "") + _beh) if (_calc or _beh) else ""
        insight(_txt + ((" &middot; Watch: " + _watch) if _watch else ""))

    # ---- trend charts on the priority metrics ----
    go = _go()
    last_n = win_months
    xl = [pd.Timestamp(m).strftime("%b %y") for m in last_n]
    _reg_colors = EXEC_REG_COLORS

    sec("Momentum", f"Where the region is heading on the priority metrics - {_win_sel} to {asof_lbl}")
    e1, e2 = st.columns(2)
    with e1:
        fig = go.Figure()
        for region in EXEC_REGIONS:
            ys = [val(region, "net_adds", m) for m in last_n]
            fig.add_trace(go.Bar(x=xl, y=ys, name=EXEC_RDISP.get(region, region),
                                 marker_color=_reg_colors.get(region, SLATE),
                                 customdata=[fmt(v, "num") for v in ys],
                                 hovertemplate="<b>%{fullData.name}</b>: %{customdata}<extra></extra>"))
        fig.update_layout(barmode="group")
        _base_layout(fig, "Net Adds by market")
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with e2:
        fig = go.Figure()
        for region in EXEC_REGIONS:
            ys = [val(region, "nrra", m) for m in last_n]
            fig.add_trace(go.Scatter(x=xl, y=ys, mode="lines+markers",
                                     name=EXEC_RDISP.get(region, region),
                                     line=dict(color=_reg_colors.get(region, SLATE), width=2.4),
                                     marker=dict(size=7, color="white",
                                                 line=dict(color=_reg_colors.get(region, SLATE), width=2)),
                                     customdata=[fmt(v, "pct") for v in ys],
                                     hovertemplate="<b>%{fullData.name}</b>: %{customdata}<extra></extra>"))
        pct_axis(_base_layout(fig, "NRRA % by market"))
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

    e3, e4 = st.columns(2)
    with e3:
        fig = go.Figure()
        for region in EXEC_REGIONS:
            ys = [val(region, "live_sold_rate", m) for m in last_n]
            fig.add_trace(go.Scatter(x=xl, y=ys, mode="lines+markers",
                                     name=EXEC_RDISP.get(region, region),
                                     line=dict(color=_reg_colors.get(region, SLATE), width=2.4),
                                     marker=dict(size=7, color="white",
                                                 line=dict(color=_reg_colors.get(region, SLATE), width=2)),
                                     customdata=[fmt(v, "pct") for v in ys],
                                     hovertemplate="<b>%{fullData.name}</b>: %{customdata}<extra></extra>"))
        pct_axis(_base_layout(fig, "Live Sold Rate by market"))
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with e4:
        # Clean filled map - no pins: countries shaded by the chosen priority metric.
        _flat = [(lbl, cn, kind) for _, items in EXEC_GROUPS for (lbl, cn, kind, _) in items
                 if cn in df.columns]
        _map_lbls = [lbl for lbl, _, _ in _flat]
        _msel = st.selectbox("Map metric", _map_lbls,
                             index=_map_lbls.index("Net Adds") if "Net Adds" in _map_lbls else 0,
                             key="exec_map_metric")
        _mcol, _mkind = next((cn, kind) for lbl, cn, kind in _flat if lbl == _msel)
        snapc = df[(df["month_end"] == asof) & (df["country"].isin(COUNTRY_ISO))].dropna(subset=[_mcol])
        if snapc.empty:
            st.info("No data to map.")
        else:
            cts = snapc["country"].tolist()
            vals = snapc[_mcol].astype(float).tolist()
            fig = gulf_map(cts, vals, _mkind,
                           f"{_msel} - {asof_lbl} (click a country for its deep dive)", height=380)
            _ev = _chart_with_select(fig, "exec_map_sel",
                                     {"displayModeBar": False, "scrollZoom": True})
            _ck = _clicked_country(_ev) if _ev is not None else None
            if _ck:
                st.session_state["_pending_cty"] = _ck
                st.session_state["_pending_view"] = "Market deep dive"
                st.rerun()


def _render_panel_overview(df, all_months, all_labels, cur_month_start):
    """Existing 'Panel Overview' - filter bar, KPI gauges, KPI cards, Revenue,
    Sales & Churn, Mix, Countries map + comparison, footer. Rendered inside the
    'Panel Overview' tab. Banner and BQ load happen once in main() and are shared
    across tabs, so both are stripped from here."""
    # ---- view toggle: Executive (all markets, the slides' mental model) vs GM deep dive ----
    # Apply a pending view jump (exec-map country click) BEFORE the widget is instantiated.
    _pv = st.session_state.pop("_pending_view", None)
    if _pv:
        st.session_state["flt_view"] = _pv
    try:
        view = st.segmented_control("View", ["Executive", "Market deep dive"],
                                    default="Executive", key="flt_view") or "Executive"
    except Exception:
        view = st.radio("View", ["Executive", "Market deep dive"], horizontal=True, key="flt_view")
    if view == "Executive":
        _render_executive_view(df, all_months, cur_month_start)
        return

    # ---- top filter bar (pills + period presets); no page title - the browser tab carries the name ----
    FLAG = {"Middle East": "\U0001F30D", "Saudi Arabia": "\U0001F1F8\U0001F1E6",
            "UAE": "\U0001F1E6\U0001F1EA", "Kuwait": "\U0001F1F0\U0001F1FC",
            "Bahrain": "\U0001F1E7\U0001F1ED", "Qatar": "\U0001F1F6\U0001F1E6"}
    countries = [c for c in COUNTRIES if c in set(df["country"])]
    disp = {c: (FLAG.get(c, "") + " " + c).strip() for c in countries}
    rev = {v: k for k, v in disp.items()}

    # Click-to-filter: apply a pending country (from map / bar clicks) BEFORE the Market
    # picker widget is instantiated (Streamlit forbids changing a widget's state after).
    _pending = st.session_state.pop("_pending_cty", None)
    if _pending and _pending in disp:
        st.session_state["flt_cty"] = disp[_pending]

    def _picker(label, options, default, key):
        """Segmented pills with graceful fallback for older Streamlit versions."""
        try:
            v = st.segmented_control(label, options, default=default, key=key)
        except Exception:
            try:
                v = st.pills(label, options, default=default, key=key)
            except Exception:
                v = st.selectbox(label, options, index=options.index(default), key=key)
        return v or default

    with st.container(border=True):
        r1a, r1b = st.columns([5.2, 0.8], vertical_alignment="bottom")
        with r1a:
            sel_disp = _picker("Market", list(disp.values()), disp[countries[0]], "flt_cty")
            sel_country = rev.get(sel_disp, countries[0])
        with r1b:
            # Refresh is a developer-only control (cache-buster); viewers get the 15-min cache.
            if _is_developer():
                try:
                    _refresh = st.button("Refresh", icon=":material/refresh:", use_container_width=True)
                except TypeError:
                    _refresh = st.button("Refresh", use_container_width=True)
                if _refresh:
                    load_bridge.clear()
                    st.rerun()
        period = _picker("Period", ["3M", "6M", "12M", "YTD", "All", "Custom"], "12M", "flt_period")
        custom_rng = None
        if period == "Custom":
            _min_d = pd.Timestamp(all_months[0]).to_pydatetime().date().replace(day=1)
            _max_d = pd.Timestamp(all_months[-1]).to_pydatetime().date()
            try:
                _pop = st.popover("Pick dates")
            except Exception:
                _pop = st.container()
            with _pop:
                _cal = st.date_input("Custom range", value=(_min_d, _max_d),
                                     min_value=_min_d, max_value=_max_d, key="flt_custom_cal")
            if isinstance(_cal, (list, tuple)) and len(_cal) == 2 and all(_cal):
                custom_rng = _cal

    # Resolve the month window from the preset.
    if period == "Custom" and custom_rng:
        _c0 = pd.Timestamp(custom_rng[0]).replace(day=1)
        _c1 = pd.Timestamp(custom_rng[1]) + pd.offsets.MonthEnd(0)
        keep = [m for m in all_months if _c0 <= pd.Timestamp(m) <= _c1] or all_months
    elif period in ("3M", "6M", "12M"):
        n = int(period[:-1])
        keep = all_months[-(n + 1):]  # last N closed months + the current partial one
    elif period == "YTD":
        _yr = pd.Timestamp.today().year
        keep = [m for m in all_months if pd.Timestamp(m).year == _yr] or all_months
    else:
        keep = all_months
    keep_months = set(keep)

    d = df[(df["country"] == sel_country) & (df["month_end"].isin(keep_months))]
    d = d.sort_values("month_end").reset_index(drop=True)
    if d.empty:
        st.info("No rows for this selection.")
        st.stop()
    x = d["month_label"].tolist()

    closed_mask = d["month_end"] < cur_month_start
    closed_idx = int(closed_mask[closed_mask].index.max()) if closed_mask.any() else len(d) - 1
    kpi_row = d.iloc[closed_idx]
    # Delta base = the FIRST closed month in the selected window, so KPI deltas answer
    # "how did we move over the selected period" and visibly react to the timeline.
    # With a single closed month in view, fall back to plain month-over-month.
    _first_closed = int(closed_mask[closed_mask].index.min()) if closed_mask.any() else 0
    if _first_closed < closed_idx:
        base_row = d.iloc[_first_closed]
        delta_label = "vs " + str(base_row["month_label"])
    elif closed_idx > 0:
        base_row = d.iloc[closed_idx - 1]
        delta_label = "vs " + str(base_row["month_label"])
    else:
        base_row = None
        delta_label = "vs prior"

    st.markdown(
        '<div class="hdr">'
        f'<span class="badge">{sel_country}</span>'
        f'<span class="badge">{x[0]} - {x[-1]} &middot; {len(d)} months</span>'
        f'<span class="badge">KPIs as of {kpi_row["month_label"]} &middot; deltas {delta_label}</span></div>',
        unsafe_allow_html=True,
    )

    def kd(col_name):
        try:
            if base_row is None or col_name not in d.columns:
                return None
            return float(kpi_row[col_name]) - float(base_row[col_name])
        except Exception:
            return None

    def kseries(col_name):
        if col_name not in d.columns:
            return None
        return d[col_name].iloc[: closed_idx + 1].fillna(0).tolist()

    # ---- KPI ring gauges (percent metrics) ----
    conc_now = None
    conc_prev = None
    if {"gross_rr_usd", "rr_after_mko_mfo_usd"} <= set(d.columns):
        try:
            g, n = float(kpi_row["gross_rr_usd"]), float(kpi_row["rr_after_mko_mfo_usd"])
            conc_now = (g - n) / g if g else None
            if base_row is not None:
                gp, np_ = float(base_row["gross_rr_usd"]), float(base_row["rr_after_mko_mfo_usd"])
                conc_prev = (gp - np_) / gp if gp else None
        except Exception:
            pass

    GAUGES = [
        ("Occupancy", kpi_row.get("occupancy"), kd("occupancy"), NAVY, True, "occupancy"),
        ("Live Sold Rate", kpi_row.get("live_sold_rate"), kd("live_sold_rate"), TEAL, True,
         "live_sold_rate"),
        ("Live Sold Rate w/ Approved", kpi_row.get("live_sold_rate_approved"),
         kd("live_sold_rate_approved"), ORANGE, True, "live_sold_rate_approved"),
        ("Concession Load", conc_now,
         (conc_now - conc_prev) if (conc_now is not None and conc_prev is not None) else None,
         RED, False, "__concession__"),
    ]
    gcols = st.columns(4)
    for (lbl, val, dv, color, up_good, gcol_name), gcol in zip(GAUGES, gcols):
        with gcol:
            with st.container(border=True):
                st.plotly_chart(donut(val, color), use_container_width=True,
                                config={"displayModeBar": False, "staticPlot": True})
                st.markdown(f'<p class="g-lbl">{lbl}</p>', unsafe_allow_html=True)
                if dv is None:
                    st.markdown('<p class="g-dlt kpi-d-na">-</p>', unsafe_allow_html=True)
                else:
                    good = (dv >= 0) if up_good else (dv <= 0)
                    cls = "kpi-d-up" if good else "kpi-d-dn"
                    arrow = "&#9650;" if dv >= 0 else "&#9660;"
                    st.markdown(f'<p class="g-dlt"><span class="{cls}">{arrow} {dv * 100:+.1f} pp</span></p>',
                                unsafe_allow_html=True)
                # "How is this calculated / behaving" card
                try:
                    with st.popover("calc & trend", use_container_width=True):
                        if gcol_name == "__concession__":
                            st.markdown("Concession Load = 1 - (RR after MKO/MFO / Gross RR): the "
                                        "share of gross license fees given away as MKO/MFO, custom, "
                                        "promo and term discounts (from the SF revenue schedules).")
                        else:
                            _calc, _watch = METRIC_INFO.get(gcol_name, ("", ""))
                            if _calc:
                                st.markdown(_calc)
                            _beh = _behavior_html(d, gcol_name, "pct", closed_idx)
                            if _beh:
                                st.markdown('<div class="nm-insight">' + _beh + "</div>",
                                            unsafe_allow_html=True)
                            if _watch:
                                st.caption("Watch: " + _watch)
                except Exception:
                    pass

    # ---- KPI number cards ----
    KPIS = [
        ("CWs", "cws", "num", True),
        ("Approved Deals", "approved_deals", "num", True),
        ("New Occupied Kitchens", "new_occupied_k", "num", True),
        ("NRRA $", "nrra_usd", "usdk", True),
        ("Gross RR $ (EoP)", "gross_rr_usd", "usdk", True),
    ]
    slots = st.columns(len(KPIS))
    for (label, col_name, kind, up_good), slot in zip(KPIS, slots):
        if col_name not in d.columns:
            continue
        dv = kd(col_name)
        dstr = fmt(abs(dv) if dv is not None else None, kind)
        _calc, _watch = METRIC_INFO.get(col_name, ("", ""))
        _beh = _behavior_html(d, col_name, kind, closed_idx)
        kpi_card(slot, label, fmt(kpi_row.get(col_name), kind), dv, dstr, kseries(col_name), up_good,
                 delta_label=delta_label, explain=(_calc, _watch, _beh))

    go = _go()

    # ---- Revenue ----
    def _top_market(col):
        """Country with the highest value of `col` at the KPI month (ME excluded)."""
        try:
            snapc = df[(df["month_end"] == kpi_row["month_end"]) & (df["country"] != "Middle East")]
            snapc = snapc.dropna(subset=[col])
            if snapc.empty:
                return None, None
            r0 = snapc.loc[snapc[col].astype(float).idxmax()]
            return r0["country"], float(r0[col])
        except Exception:
            return None, None

    sec("Revenue", "Recognized revenue flow and the occupied-kitchen revenue stock")
    _parts = []
    _nr, _dnr = kpi_row.get("nrra_usd"), kd("nrra_usd")
    if _nr is not None:
        _s = f"NRRA <b>{fmt(_nr, 'usd')}</b> in {kpi_row['month_label']}"
        if _dnr is not None:
            _s += f" ({'+' if _dnr >= 0 else '-'}{fmt(abs(_dnr), 'usd')} {delta_label})"
        _parts.append(_s)
    _tc, _tv = _top_market("nrra_usd")
    if _tc:
        _parts.append(f"top market: <b>{_tc}</b> ({fmt(_tv, 'usd')})")
    if conc_now is not None:
        _parts.append(f"concession load <b>{conc_now * 100:.1f}%</b> of gross RR")
    if _parts:
        insight(" &middot; ".join(_parts))
    c1, c2 = st.columns(2)
    with c1:
        fig = go.Figure()
        for lbl, cn, color in [("RRA $", "rra_usd", TEAL), ("RRL $", "rrl_usd", RED),
                               ("NRRA $", "nrra_usd", NAVY)]:
            if cn in d.columns:
                add_line(fig, x, d[cn].tolist(), lbl, color, closed_idx, fmt_kind="usd")
        money_axis(_base_layout(fig, "Recurring Revenue - added / lost / net"))
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with c2:
        fig = go.Figure()
        for lbl, cn, color in [("Gross RR $", "gross_rr_usd", NAVY),
                               ("RR after MKO/MFO $", "rr_after_mko_mfo_usd", ORANGE)]:
            if cn in d.columns:
                add_line(fig, x, d[cn].tolist(), lbl, color, closed_idx, fmt_kind="usd")
        money_axis(_base_layout(fig, "Occupied-kitchen revenue at EoP (gap = concessions)"))
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    _explain_block("How these revenue numbers are calculated & behaving",
                   [("RRA $ (recognized)", "rra_usd", "usd"), ("RRL $ (recognized)", "rrl_usd", "usd"),
                    ("NRRA $", "nrra_usd", "usd"), ("Gross Recurring Revenue $", "gross_rr_usd", "usd"),
                    ("RR after MKO/MFO $", "rr_after_mko_mfo_usd", "usd")], d, closed_idx)

    # ---- Sales & Churn ----
    sec("Sales &amp; Churn", "Deal flow vs attrition, and the contract value behind it")
    _parts = []
    _cw, _dcw = kpi_row.get("cws"), kd("cws")
    if _cw is not None:
        _s = f"<b>{fmt(_cw, 'num')}</b> CWs in {kpi_row['month_label']}"
        if _dcw is not None:
            _s += f" ({'+' if _dcw >= 0 else ''}{fmt(_dcw, 'num')} {delta_label})"
        _parts.append(_s)
    _tc, _tv = _top_market("cws")
    if _tc:
        _parts.append(f"most CWs: <b>{_tc}</b> ({fmt(_tv, 'num')})")
    _ch = kpi_row.get("churns_excl_transfers")
    _na = kpi_row.get("net_adds")
    if _ch is not None and _na is not None:
        _parts.append(f"churns {fmt(_ch, 'num')} &rarr; net adds <b>{fmt(_na, 'num')}</b>")
    if _parts:
        insight(" &middot; ".join(_parts))
    c3, c4 = st.columns(2)
    with c3:
        fig = go.Figure()
        if "cws" in d.columns:
            add_bar(fig, x, d["cws"].tolist(), "CWs", TEAL, closed_idx, fmt_kind="num")
        if "approved_deals" in d.columns:
            add_bar(fig, x, d["approved_deals"].tolist(), "Approved", NAVY, closed_idx, fmt_kind="num")
        if "churns_excl_transfers" in d.columns:
            add_bar(fig, x, d["churns_excl_transfers"].tolist(), "Churns", RED, closed_idx, fmt_kind="num")
        _base_layout(fig, "CWs vs Approved vs Churns")
        fig.update_layout(barmode="group")
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with c4:
        fig = go.Figure()
        if "tcv_usd" in d.columns:
            add_line(fig, x, d["tcv_usd"].tolist(), "TCV $", NAVY, closed_idx, fmt_kind="usd")
        if "approved_tcv_usd" in d.columns:
            add_line(fig, x, d["approved_tcv_usd"].tolist(), "Approved TCV $", ORANGE, closed_idx,
                     fmt_kind="usd")
        money_axis(_base_layout(fig, "TCV vs Approved TCV"))
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    _explain_block("How these sales numbers are calculated & behaving",
                   [("CWs", "cws", "num"), ("Approved Deals", "approved_deals", "num"),
                    ("Churns (excl. transfers)", "churns_excl_transfers", "num"),
                    ("Net Adds", "net_adds", "num"), ("TCV $", "tcv_usd", "usd"),
                    ("Approved TCV $", "approved_tcv_usd", "usd")], d, closed_idx)

    # ---- Mix (stacked distributions) ----
    sec("Mix", "Who the revenue comes from and how long the contracts run")
    _parts = []
    _segs = [("Start-ups", "rr_pct_startups"), ("Independents", "rr_pct_independents"),
             ("Growth", "rr_pct_growth"), ("Enterprise", "rr_pct_enterprise")]
    try:
        _seg_vals = [(lbl, float(kpi_row[cn])) for lbl, cn in _segs
                     if cn in d.columns and pd.notna(kpi_row.get(cn))]
        if _seg_vals:
            _bl, _bv = max(_seg_vals, key=lambda t: t[1])
            _parts.append(f"<b>{_bl}</b> hold {_bv * 100:.0f}% of recurring revenue")
    except Exception:
        pass
    try:
        _terms = [("&lt;= 6m", "cw_term_lte_6m"), ("7-12m", "cw_term_7_12m"),
                  ("13-18m", "cw_term_13_18m"), ("19-24m", "cw_term_19_24m"),
                  ("25-36m", "cw_term_25_36m"), ("&gt; 36m", "cw_term_gt_36m")]
        _term_vals = [(lbl, float(kpi_row[cn])) for lbl, cn in _terms
                      if cn in d.columns and pd.notna(kpi_row.get(cn))]
        if _term_vals:
            _tl, _tvv = max(_term_vals, key=lambda t: t[1])
            _parts.append(f"most common CW term: <b>{_tl}</b> ({_tvv * 100:.0f}% of wins)")
    except Exception:
        pass
    if _parts:
        insight(" &middot; ".join(_parts) + f" &middot; as of {kpi_row['month_label']}")
    c5, c6 = st.columns(2)
    with c5:
        st.plotly_chart(stacked_pct(
            d, x,
            [("Start-ups", "rr_pct_startups"), ("Independents", "rr_pct_independents"),
             ("Growth", "rr_pct_growth"), ("Enterprise", "rr_pct_enterprise")],
            "Recurring Revenue by account type", closed_idx),
            use_container_width=True, config={"displayModeBar": False})
    with c6:
        st.plotly_chart(stacked_pct(
            d, x,
            [("<= 6m", "cw_term_lte_6m"), ("7-12m", "cw_term_7_12m"),
             ("13-18m", "cw_term_13_18m"), ("19-24m", "cw_term_19_24m"),
             ("25-36m", "cw_term_25_36m"), ("> 36m", "cw_term_gt_36m")],
            "CW term mix", closed_idx),
            use_container_width=True, config={"displayModeBar": False})

    # ---- Countries (map + comparison) ----
    sec("Countries", "Click a bubble on the map or a bar to focus the whole dashboard on that market")
    labels = [l for (l, c, k) in COMPARE_METRICS if c in df.columns]
    default_lbl = "NRRA $ (net)" if "NRRA $ (net)" in labels else labels[0]
    sel_label = st.selectbox("Metric", labels, index=labels.index(default_lbl))
    sel_col, sel_kind = next((c, k) for (l, c, k) in COMPARE_METRICS if l == sel_label)

    dc = df[(df["month_end"].isin(keep_months)) & (df["country"] != "Middle East")].copy()
    dme = df[(df["month_end"].isin(keep_months)) & (df["country"] == "Middle East")].sort_values("month_end")

    # Leader / laggard takeaway for the chosen metric.
    try:
        _snapc = dc[dc["month_end"] == kpi_row["month_end"]].dropna(subset=[sel_col])
        if not _snapc.empty:
            _ld = _snapc.loc[_snapc[sel_col].astype(float).idxmax()]
            _lg = _snapc.loc[_snapc[sel_col].astype(float).idxmin()]
            _parts = [f"<b>{_ld['country']}</b> leads {sel_label} with {fmt(_ld[sel_col], sel_kind)}",
                      f"lowest: <b>{_lg['country']}</b> ({fmt(_lg[sel_col], sel_kind)})"]
            if sel_kind != "pct" and not dme.empty:
                _mesnap = dme[dme["month_end"] == kpi_row["month_end"]]
                if not _mesnap.empty and pd.notna(_mesnap[sel_col].iloc[0]) and float(_mesnap[sel_col].iloc[0]):
                    _parts.append(f"{_ld['country']} = "
                                  f"<b>{float(_ld[sel_col]) / float(_mesnap[sel_col].iloc[0]) * 100:.0f}%</b> of ME")
            insight(" &middot; ".join(_parts) + f" &middot; {kpi_row['month_label']}")
    except Exception:
        pass
    _explain_block(f"How is '{sel_label}' calculated & behaving?",
                   [(sel_label, sel_col, sel_kind)], d, closed_idx)

    m1, m2 = st.columns([3, 2])
    with m1:
        # Real Esri basemap (ArcGIS tiles - same map service the Talabat site uses), with a
        # month slider. Static per month: sturdier than plotly's frame animation on tile maps.
        try:
            _win_labels = [pd.Timestamp(m).strftime("%b %Y") for m in sorted(keep_months)]
            if len(_win_labels) > 1:
                _map_month = st.select_slider("Map month", options=_win_labels,
                                              value=(kpi_row["month_label"]
                                                     if kpi_row["month_label"] in _win_labels
                                                     else _win_labels[-1]),
                                              key="map_month")
            else:
                _map_month = _win_labels[0]
            snapm = dc[dc["month_label"] == _map_month].dropna(subset=[sel_col]).copy()
            if snapm.empty:
                st.info("No data for this month.")
            else:
                cts = snapm["country"].tolist()
                vals = snapm[sel_col].astype(float).tolist()
                fig = gulf_map(cts, vals, sel_kind,
                               f"{sel_label} - {_map_month} (click a country to focus)")
                _ev = _chart_with_select(fig, "map_sel",
                                         {"displayModeBar": False, "scrollZoom": True})
                _ck = _clicked_country(_ev) if _ev is not None else None
                if _ck and _ck != sel_country:
                    st.session_state["_pending_cty"] = _ck
                    st.rerun()
        except Exception as _map_err:
            st.info("Map unavailable: " + str(_map_err)[:160])
    with m2:
        snap = dc[dc["month_end"] == kpi_row["month_end"]].dropna(subset=[sel_col]).sort_values(sel_col)
        fig = go.Figure(go.Bar(
            x=snap[sel_col], y=snap["country"], orientation="h",
            marker_color=[COUNTRY_COLORS.get(c, SLATE) for c in snap["country"]],
            text=[fmt(v, sel_kind) for v in snap[sel_col]], textposition="outside",
            customdata=snap["country"].tolist(),
            hovertemplate="<b>%{customdata}</b>: %{text}<extra></extra>",
        ))
        _base_layout(fig, f"{sel_label} - {kpi_row['month_label']}", height=470)
        fig.update_xaxes(showticklabels=False)
        fig.update_layout(margin=dict(l=8, r=70, t=44, b=4))
        _ev2 = _chart_with_select(fig, "rank_sel", {"displayModeBar": False})
        _ck2 = _clicked_country(_ev2) if _ev2 is not None else None
        if _ck2 and _ck2 != sel_country:
            st.session_state["_pending_cty"] = _ck2
            st.rerun()

    t1, t2 = st.columns(2)
    with t1:
        fig = go.Figure()
        for cty in [c for c in COUNTRIES if c != "Middle East"]:
            dd = dc[dc["country"] == cty].sort_values("month_end")
            if not dd.empty and sel_col in dd.columns:
                fig.add_trace(go.Scatter(
                    x=dd["month_label"], y=dd[sel_col], mode="lines+markers", name=cty,
                    line=dict(color=COUNTRY_COLORS.get(cty, SLATE), width=2.2),
                    marker=dict(size=6, color="white",
                                line=dict(color=COUNTRY_COLORS.get(cty, SLATE), width=2)),
                ))
        _base_layout(fig, f"{sel_label} - trend by country", height=380)
        if sel_kind == "usd":
            money_axis(fig)
        elif sel_kind == "pct":
            pct_axis(fig)
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with t2:
        # Country share of ME, stacked by month (last 6 closed months) - only for additive metrics.
        if sel_kind != "pct":
            closed = dc[dc["month_end"] < cur_month_start].sort_values("month_end")
            last_m = sorted(closed["month_end"].unique())[-6:]
            fig = go.Figure()
            for cty in [c for c in COUNTRIES if c != "Middle East"]:
                dd = closed[(closed["country"] == cty) & (closed["month_end"].isin(last_m))]
                dd = dd.sort_values("month_end")
                if not dd.empty and sel_col in dd.columns:
                    fig.add_trace(go.Bar(
                        y=[pd.Timestamp(m).strftime("%b %Y") for m in dd["month_end"]],
                        x=dd[sel_col], name=cty, orientation="h",
                        marker_color=COUNTRY_COLORS.get(cty, SLATE),
                    ))
            fig.update_layout(barmode="stack")
            _base_layout(fig, f"{sel_label} - country mix, last 6 closed months", height=380)
            if sel_kind == "usd":
                fig.update_xaxes(tickprefix="$", tickformat="~s")
            st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
        else:
            st.caption("Country-mix stack is shown for additive metrics (counts and $), not rates.")

    # ---- branded footer ----
    fcol1, fcol2 = st.columns([0.25, 9.75], vertical_alignment="center")
    with fcol1:
        if os.path.exists(LOGO_PATH):
            st.image(LOGO_PATH, width=36)
    with fcol2:
        st.caption("**NAMAA - ME RevOps**")


# ---------------------------------------------------------------- AI analyst (NAMAA Copilot)
# Metric ids the copilot can talk about: (label, bridge column, kind, question aliases).
AI_METRICS = [
    ("Kitchens sold (CWs)", "cws", "num",
     ["kitchens sold", "closed won", "closed wons", " cws", " cw ", "sold kitchens"]),
    ("Approved Deals", "approved_deals", "num", ["approved", "approvals", "pipeline"]),
    ("Kitchens churned", "churns_excl_transfers", "num", ["churn", "churned", "churns"]),
    ("Net Adds", "net_adds", "num", ["net adds", "net add", "net growth"]),
    ("RRA %", "rra", "pct", [" rra"]),
    ("RRL %", "rrl", "pct", [" rrl", "revenue lost"]),
    ("NRRA %", "nrra", "pct", ["nrra", "net revenue"]),
    ("AE Productivity (CWs/AE)", "sales_team_cw_productivity", "num1",
     ["productivity", "cws/ae", "per ae"]),
    ("Owned sites", "all_facilities", "num", ["sites", "facilities", "owned"]),
    ("Live kitchens", "total_kitchens", "num", ["live kitchens", "total kitchens", "capacity"]),
    ("Occupied kitchens", "occupied_kitchens", "num", ["occupied"]),
    ("Live Sold %", "live_sold_rate", "pct", ["live sold", "sold rate"]),
    ("Occupancy %", "occupancy", "pct", ["occupancy", " occ "]),
    ("Gross RR $", "gross_rr_usd", "usd", ["gross rr", "gross recurring", "rr $"]),
    ("RR after MKO/MFO $", "rr_after_mko_mfo_usd", "usd",
     ["after mko", "mko", "mfo", "net rr", "net fee", "concession"]),
]
AI_COUNTRY_ALIASES = {
    "Saudi Arabia": ["saudi", "ksa", "riyadh", "jeddah"],
    "UAE": ["uae", "emirates", "dubai", "abu dhabi"],
    "Kuwait": ["kuwait"],
    "Bahrain": ["bahrain"],
    "Qatar": ["qatar", "doha"],
    "Middle East": ["middle east", "region", "overall", "whole business"],
}


def _ai_key():
    try:
        return (st.secrets.get("ANTHROPIC_API_KEY", "") or "").strip()
    except Exception:
        return ""


def _ai_model():
    try:
        return st.secrets.get("AI_MODEL", "claude-sonnet-4-6") or "claude-sonnet-4-6"
    except Exception:
        return "claude-sonnet-4-6"


@st.cache_data(ttl=900, show_spinner=False)
def _ai_digest(df):
    """Compact CSV of the whole bridge (per country x month) - the copilot's dataset."""
    cols = [(cn, kind) for (_l, cn, kind, _a) in AI_METRICS if cn in df.columns]
    lines = ["month,country," + ",".join(cn for cn, _ in cols)]
    for _, r in df.sort_values(["country", "month_end"]).iterrows():
        vals = []
        for cn, kind in cols:
            try:
                v = float(r.get(cn))
                if pd.isna(v):
                    vals.append("")
                elif kind == "pct":
                    vals.append(f"{v:.4f}")
                elif kind == "usd":
                    vals.append(f"{v:.0f}")
                else:
                    vals.append(f"{v:g}")
            except Exception:
                vals.append("")
        lines.append(pd.Timestamp(r["month_end"]).strftime("%Y-%m") + ","
                     + str(r["country"]) + "," + ",".join(vals))
    return "\n".join(lines)


def _ai_glossary():
    import re as _re
    out = []
    for lbl, cn, kind, _a in AI_METRICS:
        unit = "ratio 0-1" if kind == "pct" else ("USD" if kind == "usd" else "count")
        calc, _w = METRIC_INFO.get(cn, ("", ""))
        ln = f"- {cn} ({lbl}, {unit})"
        if calc:
            ln += ": " + _re.sub(r"<[^>]+>", "", calc)
        out.append(ln)
    return "\n".join(out)


def _ai_system(df, cur_month_start):
    closed = sorted(m for m in df["month_end"].unique() if pd.Timestamp(m) < cur_month_start)
    last_lbl = pd.Timestamp(closed[-1]).strftime("%b %Y") if closed else "n/a"
    return (
        "You are NAMAA Copilot, the AI analyst inside NAMAA's ME Sales Panel (CloudKitchens Middle "
        "East: UAE, Saudi Arabia, Kuwait, Bahrain, Qatar, plus the Middle East rollup). The last "
        f"closed month is {last_lbl}; later months are partial and must be flagged as such. Answer "
        "ONLY from the CSV dataset below - never invent numbers. Percentages are ratios (0.62 = "
        "62%). Currency is USD. Be crisp and executive: lead with the answer, then 2-4 supporting "
        "facts naming months. Speak the house language (CWs, RRA, NRRA, Live Sold).\n\n"
        "METRIC GLOSSARY\n" + _ai_glossary() + "\n\n"
        "CHART PROTOCOL: when a visual helps, append at the VERY END exactly one fenced block:\n"
        "```chart\n"
        '{"metrics": ["occupancy"], "countries": ["Saudi Arabia", "UAE"], "months": 12, '
        '"type": "line"}\n'
        "```\n"
        "metrics = up to 2 column ids from the glossary; countries from: Middle East, UAE, Saudi "
        "Arabia, Kuwait, Bahrain, Qatar; type: line or bar; months: trailing closed-month window "
        "(3-36). The app renders it - never mention the block itself.\n\n"
        "DATASET (monthly, per country)\n" + _ai_digest(df)
    )


def _claude_stream(messages, system):
    """Stream text deltas from the Anthropic Messages API over SSE (requests only, no SDK)."""
    import json as _json

    import requests as _rq

    r = _rq.post(
        "https://api.anthropic.com/v1/messages",
        headers={"x-api-key": _ai_key(), "anthropic-version": "2023-06-01",
                 "content-type": "application/json"},
        json={"model": _ai_model(), "max_tokens": 1600, "stream": True,
              "system": system, "messages": messages},
        stream=True, timeout=120)
    if r.status_code != 200:
        raise RuntimeError(f"API {r.status_code}: {r.text[:200]}")
    for raw in r.iter_lines(decode_unicode=True):
        if not raw or not raw.startswith("data:"):
            continue
        try:
            ev = _json.loads(raw[5:].strip())
        except Exception:
            continue
        if ev.get("type") == "content_block_delta":
            d = ev.get("delta", {})
            if d.get("type") == "text_delta":
                yield d.get("text", "")


def _split_ai_answer(text):
    """Separate the visible answer from the trailing ```chart``` spec block."""
    import json as _json
    import re as _re

    specs = []
    m = _re.search(r"```chart\s*(\{.*?\})\s*```", text, _re.S)
    visible = text
    if m:
        visible = (text[:m.start()] + text[m.end():]).strip()
        try:
            specs.append(_json.loads(m.group(1)))
        except Exception:
            pass
    return visible, specs


def _render_ai_chart(spec, df, cur_month_start):
    """Draw the chart the copilot asked for - validated against real columns/countries."""
    try:
        go = _go()
        byid = {cn: (lbl, kind) for (lbl, cn, kind, _a) in AI_METRICS}
        mcols = [c for c in (spec.get("metrics") or []) if c in byid and c in df.columns][:2]
        if not mcols:
            return
        cts = [c for c in (spec.get("countries") or []) if c in EXEC_REGIONS] or ["Middle East"]
        try:
            n = max(3, min(int(spec.get("months", 12)), 36))
        except Exception:
            n = 12
        typ = str(spec.get("type", "line")).lower()
        closed = sorted(m for m in df["month_end"].unique() if pd.Timestamp(m) < cur_month_start)[-n:]
        fig = go.Figure()
        for i, cn in enumerate(mcols):
            lbl, kind = byid[cn]
            for cty in cts:
                dd = df[(df["country"] == cty) & (df["month_end"].isin(closed))].sort_values("month_end")
                if dd.empty:
                    continue
                xs = [pd.Timestamp(m).strftime("%b %y") for m in dd["month_end"]]
                ys = pd.to_numeric(dd[cn], errors="coerce")
                nm = EXEC_RDISP.get(cty, cty) + (f" - {lbl}" if len(mcols) > 1 else "")
                color = EXEC_REG_COLORS.get(cty, SLATE)
                if typ == "bar":
                    fig.add_trace(go.Bar(x=xs, y=ys, name=nm, marker_color=color))
                else:
                    fig.add_trace(go.Scatter(
                        x=xs, y=ys, mode="lines+markers", name=nm,
                        line=dict(color=color, width=2.2, dash=("dot" if i else "solid")),
                        marker=dict(size=6, color="white", line=dict(color=color, width=2))))
        title = " & ".join(byid[c][0] for c in mcols) + f" - last {len(closed)} closed months"
        _base_layout(fig, title, height=360)
        kinds = {byid[c][1] for c in mcols}
        if kinds == {"pct"}:
            pct_axis(fig)
        elif kinds <= {"usd", "usdk"}:
            money_axis(fig)
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    except Exception:
        pass


def _local_ai_answer(q, df, cur_month_start):
    """Deterministic analyst for when no API key is configured - still answers and charts."""
    ql = " " + q.lower() + " "
    hit = next(((lbl, cn, kind) for lbl, cn, kind, als in AI_METRICS
                if cn in df.columns and any(a in ql for a in als)), None)
    lbl, cn, kind = hit or ("Net Adds", "net_adds", "num")
    have = set(df["country"].unique())
    cts = [c for c, als in AI_COUNTRY_ALIASES.items() if c in have and any(a in ql for a in als)]
    closed = sorted(m for m in df["month_end"].unique() if pd.Timestamp(m) < cur_month_start)
    if not closed:
        return "No closed months in the bridge yet.", []
    la = closed[-1]
    pr = closed[-2] if len(closed) > 1 else None

    def gv(cty, m):
        r = df[(df["country"] == cty) & (df["month_end"] == m)]
        try:
            v = float(r[cn].iloc[0])
            return None if pd.isna(v) else v
        except Exception:
            return None

    show = cts or EXEC_REGIONS
    lines = [f"**{lbl}** as of **{pd.Timestamp(la).strftime('%b %Y')}** (last closed month):", ""]
    rank = []
    for cty in show:
        v = gv(cty, la)
        if v is None:
            continue
        s = f"- **{EXEC_RDISP.get(cty, cty)}**: {fmt(v, kind)}"
        if pr is not None:
            p = gv(cty, pr)
            if p is not None:
                dv = v - p
                dtxt = f"{dv * 100:+.1f} pp" if kind == "pct" else ("+" if dv >= 0 else "") + fmt(dv, kind)
                s += f" ({dtxt} MoM)"
        lines.append(s)
        if cty != "Middle East":
            rank.append((cty, v))
    if len(rank) > 1:
        b = max(rank, key=lambda t: t[1])
        w = min(rank, key=lambda t: t[1])
        lines += ["", f"Strongest: **{EXEC_RDISP.get(b[0], b[0])}** ({fmt(b[1], kind)}) - "
                      f"weakest: **{EXEC_RDISP.get(w[0], w[0])}** ({fmt(w[1], kind)})."]
    lines += ["", "_Offline analyst - add ANTHROPIC_API_KEY to secrets to unlock the full copilot._"]
    spec = {"metrics": [cn], "countries": (cts or ["Middle East", "Saudi Arabia", "UAE"])[:3],
            "months": 12, "type": "line"}
    return "\n".join(lines), [spec]


def _anomaly_scan(df, asof, months_closed):
    """Z-score scan of last MoM moves vs each series' own history - the anomaly radar."""
    out = []
    hist = [m for m in months_closed if pd.Timestamp(m) <= pd.Timestamp(asof)][-13:]
    if len(hist) < 6:
        return out
    for lbl, cn, kind, _a in AI_METRICS:
        if cn not in df.columns:
            continue
        for region in EXEC_REGIONS:
            s = df[(df["country"] == region) & (df["month_end"].isin(hist))].sort_values("month_end")
            v = pd.to_numeric(s[cn], errors="coerce")
            if v.notna().sum() < 6:
                continue
            d = v.diff().dropna()
            if len(d) < 5:
                continue
            sd = float(d.iloc[:-1].std() or 0)
            if sd == 0:
                continue
            z = (float(d.iloc[-1]) - float(d.iloc[:-1].mean())) / sd
            if abs(z) >= 2.2:
                out.append({"region": region, "metric": lbl, "kind": kind,
                            "z": z, "delta": float(d.iloc[-1])})
    out.sort(key=lambda a: -abs(a["z"]))
    return out[:6]


def _render_ai_analyst(df, all_months, cur_month_start):
    """The Copilot tab: streaming chat over the bridge + charts drawn on request."""
    key_on = bool(_ai_key())
    st.markdown(
        '<div class="ai-hero"><div class="ai-spark">&#10022;</div><div>'
        '<p class="ai-t">NAMAA Copilot</p>'
        '<p class="ai-s">Ask anything about the ME Sales Panel - it answers with the numbers '
        'and draws the chart for you.</p></div>'
        + ('<span class="ai-badge on">Claude connected</span>' if key_on else
           '<span class="ai-badge off">offline engine</span>')
        + "</div>", unsafe_allow_html=True)

    if "ai_chat" not in st.session_state:
        st.session_state["ai_chat"] = []

    suggestions = [
        "How did KSA occupancy trend over the last 12 months?",
        "Compare net adds across all markets",
        "Which market has the best AE productivity?",
        "Where are churns concentrated?",
        "How big is the concession load (Gross RR vs after MKO/MFO)?",
        "Is UAE live sold rate improving?",
    ]
    pend = st.session_state.pop("_ai_pending", None)
    scols = st.columns(3)
    for i, s in enumerate(suggestions):
        with scols[i % 3]:
            if st.button(s, key=f"ai_sug_{i}", use_container_width=True):
                pend = s

    _avatar = LOGO_PATH if os.path.exists(LOGO_PATH) else None
    for msg in st.session_state["ai_chat"]:
        with st.chat_message(msg["role"], avatar=(_avatar if msg["role"] == "assistant" else None)):
            st.markdown(msg["visible"])
            for sp in msg.get("specs", []):
                _render_ai_chart(sp, df, cur_month_start)

    q = st.chat_input("Ask the panel - e.g. 'compare UAE and KSA live sold rate'") or pend
    if not q:
        return
    st.session_state["ai_chat"].append({"role": "user", "visible": q, "raw": q, "specs": []})
    with st.chat_message("user"):
        st.markdown(q)
    with st.chat_message("assistant", avatar=_avatar):
        ph = st.empty()
        visible = specs = None
        raw = ""
        if key_on:
            try:
                msgs = [{"role": m["role"], "content": m["raw"]}
                        for m in st.session_state["ai_chat"][-9:]]
                for chunk in _claude_stream(msgs, _ai_system(df, cur_month_start)):
                    raw += chunk
                    ph.markdown(raw.split("```chart")[0] + " &#9612;", unsafe_allow_html=True)
                visible, specs = _split_ai_answer(raw)
            except Exception as e:
                st.caption("Claude call failed (" + str(e)[:120] + ") - offline answer below.")
                visible = None
        if visible is None:
            visible, specs = _local_ai_answer(q, df, cur_month_start)
            raw = visible
        ph.markdown(visible)
        for sp in specs:
            _render_ai_chart(sp, df, cur_month_start)
    st.session_state["ai_chat"].append(
        {"role": "assistant", "visible": visible, "raw": raw, "specs": specs})
    st.session_state["ai_chat"] = st.session_state["ai_chat"][-30:]


# ================================================================ Executive redesign (v2)
# Clean "boardroom" overview matching the approved concept: green header bar, a strip of
# serif-number KPI cards (one highlighted, rich hover tooltip), and per-country line charts.
# Data comes from the same me_sales_panel_k_monthly bridge; no new columns.

# (label, bridge column, format kind, up_is_good). "mo" = months suffix.
NEW_KPIS = [
    ("Occupancy",        "occupancy",        "pct",  True),
    ("Live-Sold Rate",   "live_sold_rate",   "pct",  True),
    ("Churn Rate",       "rrl",              "pct",  False),
    ("New CWs",          "cws",              "num",  True),
    ("Approved Deals",   "approved_deals",   "num",  True),
    ("Rev. Added",       "rra_usd",          "kusd", True),
    ("Net Rev. Added",   "nrra_usd",         "kusd", True),
]
NEW_HILITE = "Occupancy"
NEW_CHART_COUNTRIES = ["UAE", "Saudi Arabia", "Kuwait", "Bahrain"]

# Country snapshot table columns: (header, bridge col, format kind, colored, is_bar).
SNAP_COLS = [
    ("Occupancy",  "occupancy",             "pct",  False, True),
    ("Live-Sold",  "live_sold_rate",        "pct",  False, False),
    ("Churn Rate", "rrl",                   "pct",  "churn", False),
    ("New CWs",    "cws",                   "num",  False, False),
    ("Approved",   "approved_deals",        "num",  False, False),
    ("RRA",        "rra_usd",               "kusd", False, False),
    ("RRL",        "rrl_usd",               "kusd", "loss", False),
    ("NRRA",       "nrra_usd",              "kusd", "signed", False),
]
SNAP_ROWS = ["Middle East", "UAE", "Kuwait", "Saudi Arabia", "Bahrain"]
SNAP_FLAG = {"Middle East": "\U0001F30D", "UAE": "\U0001F1E6\U0001F1EA",
             "Kuwait": "\U0001F1F0\U0001F1FC", "Saudi Arabia": "\U0001F1F8\U0001F1E6",
             "Bahrain": "\U0001F1E7\U0001F1ED", "Qatar": "\U0001F1F6\U0001F1E6"}


def _inject_exec_css():
    """One-time CSS for the v2 executive look (serif numbers, header bar, KPI cards,
    hover tooltips, terracotta tab underline)."""
    if st.session_state.get("_exec2_css"):
        return
    st.session_state["_exec2_css"] = True
    st.markdown(
        """
        <style>
        /* Rich serif for headline numbers (Fraunces), clean sans for UI (Inter).
           Both fall back gracefully to system fonts if the CDN is blocked. */
        @import url('https://fonts.googleapis.com/css2?family=Fraunces:opsz,wght@9..144,500;9..144,600;9..144,700&family=Inter:wght@500;600;700;800&display=swap');
        :root{--nm-serif:"Fraunces",Georgia,"Times New Roman",serif;
              --nm-sans:"Inter",-apple-system,BlinkMacSystemFont,"Segoe UI",system-ui,sans-serif;}

        /* terracotta active-tab underline + clean tab labels */
        .stTabs [data-baseweb="tab-list"]{gap:2px;border-bottom:1px solid #E0DCCE;}
        .stTabs [data-baseweb="tab"]{font-family:var(--nm-sans);font-weight:700;color:#7C776A;font-size:0.95rem;}
        .stTabs [aria-selected="true"]{color:#21362B;}
        .stTabs [data-baseweb="tab-highlight"]{background-color:#D97757 !important;}

        /* executive header bar */
        .xh{display:flex;align-items:center;gap:16px;background:linear-gradient(135deg,#21362B,#2C4A3B);
            border-radius:16px;padding:16px 22px;box-shadow:0 8px 30px rgba(33,54,43,.14);
            position:relative;overflow:hidden;margin-bottom:6px;}
        .xh::after{content:"";position:absolute;right:-40px;top:-70px;width:300px;height:300px;
            background:radial-gradient(circle at center,rgba(217,119,87,.30),transparent 62%);pointer-events:none;}
        .xh img{height:44px;border-radius:9px;}
        .xh-eyebrow{font-family:var(--nm-sans);color:#9DB3A6;font-size:.66rem;font-weight:800;letter-spacing:.22em;text-transform:uppercase;margin:0;}
        .xh-title{font-family:var(--nm-serif);color:#fff;font-size:1.5rem;font-weight:600;letter-spacing:-.01em;margin:2px 0 0;}
        .xh-right{margin-left:auto;z-index:1;color:#C9D5CC;font-size:.78rem;font-weight:600;
            font-family:var(--nm-sans);display:flex;align-items:center;gap:8px;}

        /* KPI card strip */
        .xk-grid{display:grid;grid-template-columns:repeat(7,1fr);gap:12px;margin:18px 0 6px;}
        .xk{position:relative;background:#FBFAF6;border:1px solid #E0DCCE;border-radius:14px;
            padding:15px 16px 13px;box-shadow:0 2px 10px rgba(33,54,43,.06);overflow:visible;}
        .xk.hi{background:#21362B;border-color:#21362B;}
        .xk.hi .xk-label{color:#9DB3A6;}
        .xk.hi .xk-num{color:#fff;}
        .xk.hi .xk-dl{color:#C9D5CC;}
        .xk-label{font-family:var(--nm-sans);font-size:.64rem;font-weight:800;letter-spacing:.08em;
            text-transform:uppercase;color:#A79E8B;margin:0;}
        .xk-num{font-family:var(--nm-serif);font-size:2.05rem;font-weight:600;
            color:#21362B;margin:6px 0 5px;line-height:1;letter-spacing:-.01em;
            font-variant-numeric:lining-nums;}
        .xk-dl{font-family:var(--nm-sans);font-size:.72rem;font-weight:600;color:#7C776A;display:flex;align-items:center;gap:5px;}
        .xk-dl .up{color:#3F7A52;font-weight:700;}.xk-dl .dn{color:#B4472E;font-weight:700;}
        .xk-spark{margin-top:11px;display:block;}
        .xk-tip{position:absolute;left:50%;bottom:calc(100% + 10px);transform:translateX(-50%);
            width:230px;background:#21362B;color:#EAF1EC;border-radius:12px;padding:12px 14px;
            font-family:var(--nm-sans);box-shadow:0 12px 34px rgba(33,54,43,.32);
            opacity:0;visibility:hidden;transition:.14s;z-index:50;}
        .xk:hover .xk-tip{opacity:1;visibility:visible;}
        .xk-tip::after{content:"";position:absolute;left:50%;top:100%;transform:translateX(-50%);
            border:7px solid transparent;border-top-color:#21362B;}
        .xk-tip-row{display:flex;justify-content:space-between;font-size:.74rem;margin-bottom:5px;}
        .xk-tip-row span{color:#9DB3A6;}.xk-tip-row b{font-weight:800;}
        .xk-tip-txt{font-size:.72rem;color:#C9D5CC;line-height:1.35;margin-top:8px;
            border-top:1px solid rgba(255,255,255,.12);padding-top:8px;}
        @media(max-width:1100px){.xk-grid{grid-template-columns:repeat(4,1fr);}}
        @media(max-width:640px){.xk-grid{grid-template-columns:repeat(2,1fr);}}

        /* Country snapshot table */
        .snap-card{background:#FBFAF6;border:1px solid #E0DCCE;border-radius:16px;
            padding:20px 22px;box-shadow:0 2px 10px rgba(33,54,43,.06);margin:14px 0 8px;overflow-x:auto;}
        .snap-h{font-family:var(--nm-serif);font-size:1.15rem;font-weight:600;color:#21362B;margin:0;}
        .snap-sub{font-family:var(--nm-sans);font-size:.78rem;color:#A79E8B;margin:2px 0 14px;}
        table.snap{width:100%;border-collapse:collapse;font-family:var(--nm-sans);min-width:820px;}
        table.snap th{text-align:right;font-size:.64rem;font-weight:800;letter-spacing:.06em;
            text-transform:uppercase;color:#A79E8B;padding:6px 12px 10px;border-bottom:1px solid #E0DCCE;}
        table.snap th:first-child,table.snap td:first-child{text-align:left;}
        table.snap td{text-align:right;padding:12px;border-bottom:1px solid #EEEBE1;
            color:#21362B;font-weight:600;font-size:.9rem;white-space:nowrap;}
        table.snap tr:last-child td{border-bottom:none;}
        table.snap tr.me td{background:#F1EFE7;font-weight:800;}
        table.snap tr.me td:first-child{border-radius:8px 0 0 8px;}
        table.snap tr.me td:last-child{border-radius:0 8px 8px 0;}
        .snap-mkt{display:flex;align-items:center;gap:9px;font-weight:700;}
        .snap-flag{font-size:1rem;}
        .snap-occ{display:flex;align-items:center;gap:10px;justify-content:flex-end;}
        .snap-bar{width:64px;height:7px;border-radius:4px;background:#E6E2D6;overflow:hidden;flex:0 0 auto;}
        .snap-bar>i{display:block;height:100%;border-radius:4px;}
        .snap .pos{color:#3F7A52;}.snap .neg{color:#B4472E;}.snap .warn{color:#B4472E;}
        </style>
        """,
        unsafe_allow_html=True,
    )


def _kfmt(v, kind):
    if v is None:
        return "—"
    try:
        v = float(v)
    except Exception:
        return "—"
    if math.isnan(v):
        return "—"
    if kind == "pct":
        return f"{v * 100:.1f}%"
    if kind == "mo":
        return f"{v:.0f}mo"
    if kind == "num":
        return f"{v:,.0f}"
    if kind == "kusd":
        # thousands, no suffix — matches the All Hands slide style ("$184")
        neg = "-" if v < 0 else ""
        return f"{neg}${abs(v) / 1000:,.0f}"
    return fmt(v, kind)


def _kdelta_txt(v, kind):
    """Absolute-delta text (used in hover tooltip)."""
    if v is None:
        return "—"
    if kind == "pct":
        return f"{v * 100:+.2f}pp"
    if kind == "mo":
        return f"{v:+.1f}mo"
    if kind == "kusd":
        return ("+" if v >= 0 else "-") + f"${abs(v) / 1000:,.0f}"
    return f"{v:+,.0f}"


def _face_delta(cur, base, kind, up_good):
    """Card-face delta shown under the number: percentage-points for rates,
    relative % change for counts & revenue (matches the concept). Returns
    (text, css_class, arrow_html)."""
    if cur is None or base is None:
        return "—", "", ""
    if kind == "pct":
        d = (cur - base) * 100.0
        txt = f"{d:+.2f}pp"
        pos = d >= 0
    else:
        if base == 0:
            return "n/a", "", ""
        d = (cur - base) / abs(base) * 100.0
        txt = f"{d:+.1f}%"
        pos = d >= 0
    good = pos if up_good else (not pos)
    cls = "up" if good else "dn"
    ar = "&#9650;" if pos else "&#9660;"
    return txt, cls, ar


def _render_country_snapshot(df, asof, *, title="Country Snapshot", cols=None, rows=None):
    """All-Hands-style table: key metrics per market for the as-of month, ME row highlighted."""
    cols = cols or SNAP_COLS
    rows = rows or SNAP_ROWS
    asof_lbl = pd.Timestamp(asof).strftime("%b %Y")

    def gv(region, col):
        r = df[(df["country"] == region) & (df["month_end"] == asof)]
        if r.empty or col not in r.columns:
            return None
        try:
            x = float(r[col].iloc[0])
            return None if math.isnan(x) else x
        except Exception:
            return None

    head = "".join(f"<th>{h}</th>" for (h, *_ ) in cols)
    body = []
    for region in rows:
        is_me = region == "Middle East"
        disp = region  # full market names in the snapshot, like the concept
        flag = SNAP_FLAG.get(region, "")
        tds = [f'<td><span class="snap-mkt"><span class="snap-flag">{flag}</span>{disp}</span></td>']
        for (_h, col, kind, colored, is_bar) in cols:
            v = gv(region, col)
            txt = _kfmt(v, kind)
            if is_bar and v is not None:
                pct = max(0.0, min(1.0, v)) * 100.0
                barcol = COUNTRY_COLORS.get(region, NAVY) if not is_me else NAVY
                cell = (f'<div class="snap-occ"><span class="snap-bar">'
                        f'<i style="width:{pct:.0f}%;background:{barcol}"></i></span>{txt}</div>')
                tds.append(f"<td>{cell}</td>")
                continue
            klass = ""
            if colored == "churn" and v is not None and v >= 0.05:
                klass = "warn"
            elif colored == "loss":
                klass = "neg"
            elif colored == "signed" and v is not None:
                klass = "pos" if v >= 0 else "neg"
                if v >= 0:
                    txt = "+" + txt
            tds.append(f'<td class="{klass}">{txt}</td>')
        body.append(f'<tr class="{"me" if is_me else ""}">' + "".join(tds) + "</tr>")

    st.markdown(
        '<div class="snap-card">'
        f'<p class="snap-h">{title} &mdash; {asof_lbl}</p>'
        '<p class="snap-sub">Key metrics by market, current month</p>'
        '<table class="snap"><thead><tr><th>Market</th>' + head + '</tr></thead>'
        '<tbody>' + "".join(body) + '</tbody></table></div>',
        unsafe_allow_html=True,
    )


def _country_line(fig_col, df, months_closed, col, title, *, countries, alert=None, is_pct=True):
    """Smooth (spline) country lines in the brand palette — the 'fixed' chart style."""
    go = _go()
    lbls = [pd.Timestamp(m).strftime("%b %y") for m in months_closed]
    fig = go.Figure()
    for cty in countries:
        ys = _exec_series(df, cty, col, months_closed)
        fig.add_trace(go.Scatter(
            x=lbls, y=ys, mode="lines", name=EXEC_RDISP.get(cty, cty),
            line=dict(color=COUNTRY_COLORS.get(cty, SLATE), width=2.6, shape="spline", smoothing=1.0),
            connectgaps=True,
            hovertemplate="<b>%{fullData.name}</b>: %{y}<extra></extra>",
        ))
    if alert is not None:
        fig.add_hline(y=alert, line=dict(color="#B4472E", width=1.3, dash="dash"),
                      annotation_text=f"{alert*100:.0f}% threshold", annotation_position="top right",
                      annotation_font=dict(size=10, color="#B4472E"))
    fig.update_layout(
        title=dict(text=title, font=dict(size=15, color="#21362B", family="Georgia, serif")),
        height=360, margin=dict(l=6, r=6, t=42, b=4),
        plot_bgcolor="rgba(0,0,0,0)", paper_bgcolor="rgba(0,0,0,0)",
        legend=dict(orientation="h", y=-0.16, font=dict(size=11, color="#7C776A")),
        hovermode="x unified", font=dict(family="Inter, Arial", size=11, color="#7C776A"),
    )
    fig.update_xaxes(showgrid=False, tickfont=dict(size=10, color="#A79E8B"))
    fig.update_yaxes(gridcolor="#E7E4D8", griddash="dot", zerolinecolor="#E7E4D8",
                     tickfont=dict(size=10, color="#A79E8B"))
    if is_pct:
        fig.update_yaxes(tickformat=".0%")
    with fig_col:
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})


def _exec_series(df, region, col, months):
    out = []
    for m in months:
        r = df[(df["country"] == region) & (df["month_end"] == m)]
        if r.empty or col not in r.columns:
            out.append(None)
            continue
        try:
            x = float(r[col].iloc[0])
            out.append(None if math.isnan(x) else x)
        except Exception:
            out.append(None)
    return out


def _spark_svg(vals, color, w=118, h=30):
    pts = [v for v in vals if v is not None]
    if len(pts) < 2:
        return f'<svg class="xk-spark" width="100%" height="{h}" viewBox="0 0 {w} {h}"></svg>'
    lo, hi = min(pts), max(pts)
    rng = (hi - lo) or 1.0
    n = len(vals)
    coords = []
    for i, v in enumerate(vals):
        if v is None:
            continue
        px = (i / (n - 1)) * w
        py = h - ((v - lo) / rng) * (h - 4) - 2
        coords.append(f"{px:.1f},{py:.1f}")
    poly = " ".join(coords)
    return (f'<svg class="xk-spark" width="100%" height="{h}" viewBox="0 0 {w} {h}" preserveAspectRatio="none">'
            f'<polyline fill="none" stroke="{color}" stroke-width="2" points="{poly}"/></svg>')


def _render_exec_header(df, all_months):
    _inject_exec_css()
    _b64 = _logo_b64()
    img = (f'<img src="data:image/jpeg;base64,{_b64}"/>' if _b64 else "")
    now = pd.Timestamp.now()
    ts = now.strftime("%b %d, %Y &middot; %I:%M %p")
    st.markdown(
        '<div class="xh">' + img +
        '<div><p class="xh-eyebrow">Executive Sales Panel</p>'
        '<p class="xh-title">Middle East &mdash; Cloud Kitchen Network</p></div>'
        '<div class="xh-right">&#8635; Data refreshed: ' + ts + ' UAE</div></div>',
        unsafe_allow_html=True,
    )


def _render_overview_tab(df, all_months, cur_month_start):
    """v2 Overview: 7 serif KPI cards (ME headline, hover MoM/YoY) + occupancy & churn by country."""
    _inject_exec_css()
    months_closed = [m for m in all_months if pd.Timestamp(m) < cur_month_start]
    if len(months_closed) < 2:
        st.info("Not enough closed months yet.")
        return
    asof = months_closed[-1]
    prev = months_closed[-2]
    asof_ts = pd.Timestamp(asof)
    yoy_ts = asof_ts - pd.DateOffset(years=1)
    yoy = next((m for m in months_closed if pd.Timestamp(m) == yoy_ts), None)
    yoy_lbl = yoy_ts.strftime("%b '%y")
    prev_lbl = pd.Timestamp(prev).strftime("%b '%y")
    win = months_closed[-13:]  # ~1yr sparkline window

    def _v(col, m):
        if m is None:
            return None
        r = df[(df["country"] == "Middle East") & (df["month_end"] == m)]
        if r.empty or col not in r.columns:
            return None
        try:
            x = float(r[col].iloc[0])
            return None if math.isnan(x) else x
        except Exception:
            return None

    cards = []
    for label, col, kind, up_good in NEW_KPIS:
        cur = _v(col, asof)
        pv = _v(col, prev)
        yv = _v(col, yoy)
        mom = (cur - pv) if (cur is not None and pv is not None) else None
        yoy_d = (cur - yv) if (cur is not None and yv is not None) else None
        # face delta = YoY, shown pp for rates / relative-% for counts & revenue
        face_txt, cls, ar = _face_delta(cur, yv, kind, up_good)
        spark = _spark_svg(_exec_series(df, "Middle East", col, win),
                           "#fff" if label == NEW_HILITE else "#5F8575")
        narr = f"{label} across the ME network as of {asof_ts.strftime('%b %Y')}. " \
               f"MoM {_kdelta_txt(mom, kind)}, YoY {_kdelta_txt(yoy_d, kind)}."
        hi = " hi" if label == NEW_HILITE else ""
        cards.append(
            f'<div class="xk{hi}">'
            f'<p class="xk-label">{label}</p>'
            f'<p class="xk-num">{_kfmt(cur, kind)}</p>'
            f'<p class="xk-dl"><span class="{cls}">{ar} {face_txt}</span>'
            f'<span style="color:#A79E8B;font-weight:600"> vs {yoy_lbl}</span></p>'
            f'{spark}'
            f'<div class="xk-tip">'
            f'<div class="xk-tip-row"><span>MoM vs {prev_lbl}</span><b>{_kdelta_txt(mom, kind)}</b></div>'
            f'<div class="xk-tip-row"><span>YoY vs {yoy_lbl}</span><b>{_kdelta_txt(yoy_d, kind)}</b></div>'
            f'<div class="xk-tip-txt">{narr}</div>'
            f'</div></div>'
        )
    st.markdown('<div class="xk-grid">' + "".join(cards) + "</div>", unsafe_allow_html=True)

    # ---- All-Hands country snapshot ----
    _render_country_snapshot(df, asof)

    # ---- smooth country charts ----
    c1, c2 = st.columns(2)
    _country_line(c1, df, months_closed, "occupancy", "Occupancy by Country",
                  countries=NEW_CHART_COUNTRIES, is_pct=True)
    _country_line(c2, df, months_closed, "rrl", "Churn Rate by Country  ·  lower is better",
                  countries=NEW_CHART_COUNTRIES[:3], alert=0.05, is_pct=True)

    st.caption(f"KPI cards are Middle East aggregate as of **{asof_ts.strftime('%b %Y')}** (last closed month); "
               f"delta is vs {yoy_lbl} (pp for rates, % for volumes). Snapshot is current month by market. "
               f"$ figures in thousands. Source: `{BRIDGE_TABLE}`.")


def _render_revenue_tab(df, all_months, cur_month_start):
    """v2 Revenue: recognized RR flow + occupied-kitchen RR stock + RRA $ by country."""
    _inject_exec_css()
    months_closed = [m for m in all_months if pd.Timestamp(m) < cur_month_start]
    if len(months_closed) < 2:
        st.info("Not enough closed months yet.")
        return
    go = _go()
    lbls = [pd.Timestamp(m).strftime("%b %y") for m in months_closed]
    closed_idx = len(lbls) - 1
    _render_country_snapshot(df, months_closed[-1], title="Revenue Snapshot", cols=[
        ("RRA", "rra_usd", "kusd", False, False),
        ("RRL", "rrl_usd", "kusd", "loss", False),
        ("NRRA", "nrra_usd", "kusd", "signed", False),
        ("Gross RR", "gross_rr_usd", "kusd", False, False),
        ("TCV", "tcv_usd", "kusd", False, False),
    ])
    sec("Recurring Revenue", "Recognized License-Fee revenue added, lost and net (Middle East)")
    c1, c2 = st.columns(2)
    with c1:
        fig = go.Figure()
        for lbl, col, color in [("RRA $ (added)", "rra_usd", TEAL),
                                ("RRL $ (lost)", "rrl_usd", RED),
                                ("NRRA $ (net)", "nrra_usd", NAVY)]:
            add_line(fig, lbls, _exec_series(df, "Middle East", col, months_closed),
                     lbl, color, closed_idx, fmt_kind="usd")
        money_axis(_base_layout(fig, "RR added / lost / net", height=340))
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with c2:
        fig = go.Figure()
        for lbl, col, color in [("Gross RR $ (EoP)", "gross_rr_usd", NAVY),
                                ("TCV $", "tcv_usd", YELLOW)]:
            add_line(fig, lbls, _exec_series(df, "Middle East", col, months_closed),
                     lbl, color, closed_idx, fmt_kind="usd")
        money_axis(_base_layout(fig, "Gross RR (stock) & TCV committed", height=340))
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

    sec("RRA $ by Country", "Recognized revenue added per market")
    fig = go.Figure()
    for cty in NEW_CHART_COUNTRIES:
        add_line(fig, lbls, _exec_series(df, cty, "rra_usd", months_closed),
                 EXEC_RDISP.get(cty, cty), COUNTRY_COLORS.get(cty, SLATE), closed_idx, fmt_kind="usd")
    money_axis(_base_layout(fig, "RRA $ — by market", height=360))
    st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})


def _render_operations_tab(df, all_months, cur_month_start):
    """v2 Operations: sales flow (CWs/Approved/Churns), net adds by country, kitchens & sites."""
    _inject_exec_css()
    months_closed = [m for m in all_months if pd.Timestamp(m) < cur_month_start]
    if len(months_closed) < 2:
        st.info("Not enough closed months yet.")
        return
    go = _go()
    lbls = [pd.Timestamp(m).strftime("%b %y") for m in months_closed]
    closed_idx = len(lbls) - 1
    _render_country_snapshot(df, months_closed[-1], title="Operations Snapshot", cols=[
        ("Occupancy", "occupancy", "pct", False, True),
        ("New CWs", "cws", "num", False, False),
        ("Approved", "approved_deals", "num", False, False),
        ("Churns", "churns_excl_transfers", "num", False, False),
        ("Net Adds", "net_adds", "num", "signed", False),
        ("Live-Sold", "live_sold_rate", "pct", False, False),
    ])
    sec("Sales & Churn", "Contract wins, approved deals and churns across the ME network")
    c1, c2 = st.columns(2)
    with c1:
        fig = go.Figure()
        for lbl, col, color in [("CWs", "cws", TEAL), ("Approved", "approved_deals", NAVY),
                                ("Churns", "churns_excl_transfers", RED)]:
            add_bar(fig, lbls, _exec_series(df, "Middle East", col, months_closed),
                    lbl, color, closed_idx, fmt_kind="num")
        _base_layout(fig, "CWs · Approved · Churns", height=340)
        fig.update_layout(barmode="group")
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with c2:
        fig = go.Figure()
        for cty in NEW_CHART_COUNTRIES:
            add_line(fig, lbls, _exec_series(df, cty, "net_adds", months_closed),
                     EXEC_RDISP.get(cty, cty), COUNTRY_COLORS.get(cty, SLATE), closed_idx, fmt_kind="num")
        _base_layout(fig, "Net Adds by Country", height=340)
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})

    sec("Capacity", "Live kitchens and owned sites (Middle East)")
    c3, c4 = st.columns(2)
    with c3:
        fig = go.Figure()
        add_line(fig, lbls, _exec_series(df, "Middle East", "total_kitchens", months_closed),
                 "Live kitchens", NAVY, closed_idx, fmt_kind="num")
        add_line(fig, lbls, _exec_series(df, "Middle East", "occupied_kitchens", months_closed),
                 "Occupied kitchens", TEAL, closed_idx, fmt_kind="num")
        _base_layout(fig, "Kitchens — live vs occupied", height=340)
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})
    with c4:
        fig = go.Figure()
        add_line(fig, lbls, _exec_series(df, "Middle East", "all_facilities", months_closed),
                 "Owned sites", ORANGE, closed_idx, fmt_kind="num")
        _base_layout(fig, "Owned sites", height=340)
        st.plotly_chart(fig, use_container_width=True, config={"displayModeBar": False})


# ---------------------------------------------------------------- main
def main():
    """Access gate + load BQ bridge once, then a NAMAA exec header and five tabs:
    Overview / Revenue / Operations / ME All Hands / CR All Hands."""
    _access_gate()

    try:
        df = load_bridge()
    except Exception as e:
        st.error("Could not load the bridge from BigQuery. Check credentials.\n\nDetails: " + str(e))
        st.stop()
    if df.empty:
        st.warning("The bridge returned no rows.")
        st.stop()

    df["month_label"] = df["month_end"].dt.strftime("%b %Y")
    all_months = sorted(df["month_end"].unique())
    all_labels = [pd.Timestamp(m).strftime("%b %Y") for m in all_months]
    cur_month_start = pd.Timestamp.today().normalize().replace(day=1)

    _render_exec_header(df, all_months)

    tab_ov, tab_rev, tab_ops, tab_ck, tab_cr = st.tabs([
        "Overview",
        "Revenue",
        "Operations",
        "ME All Hands",
        "CR All Hands",
    ])
    with tab_ov:
        _render_overview_tab(df, all_months, cur_month_start)
    with tab_rev:
        _render_revenue_tab(df, all_months, cur_month_start)
    with tab_ops:
        _render_operations_tab(df, all_months, cur_month_start)
    with tab_ck:
        _render_all_hands_scorecard(df, all_months, CK_SCORECARD_METRICS, "ck",
                                    pptx_title="ME All Hands - Cloud Kitchens")
    with tab_cr:
        _render_all_hands_scorecard(df, all_months, CR_SCORECARD_METRICS, "cr",
                                    pptx_title="ME All Hands - Cloud Retail")


main()
